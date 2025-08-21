#!/usr/bin/env python3
"""
website_to_ppt.py - Fixed Version
=================================
End‑to‑end pipeline that
 1. Crawls every reachable page on a target website using a **Browserbase**
    remote Chromium instance (with optional Stagehand AI actions).
 2. Extracts & cleans readable text from each page (Trafilatura).
 3. Streams the cleaned corpus + a fixed system prompt to **Claude 4
    Sonnet** (Anthropic SDK) and receives structured content for **five
    slides**.
 4. Inserts that content into a PowerPoint deck based on a user‑supplied
    `.pptx` template (python‑pptx).

Assumptions / prerequisites
---------------------------
* Python ≥ 3.10.
* Packages (install with `pip install -r requirements.txt`):
  playwright  browserbase-python  stagehand-python  anthropic  
  trafilatura  python-pptx  tiktoken  tenacity  tqdm
* Environment variables:
    BROWSERBASE_API_KEY   – your Browserbase API key
    BROWSERBASE_PROJECT_ID – your Browserbase project ID
    ANTHROPIC_API_KEY     – your Anthropic API key
* A template deck at TEMPLATE_PPTX (title + content placeholders).

Usage
-----
$ python website_to_ppt.py \
        --url https://example.com \
        --template my_template.pptx \
        --output  example_deck.pptx \
        [--stagehand]
"""

from __future__ import annotations

import argparse
import asyncio
import json
import os
import pathlib
import re
import sys
import time
from dataclasses import dataclass, field
from typing import Dict, List, Set, Tuple
from urllib.parse import urljoin, urlparse

import anthropic
import tiktoken
import trafilatura
from playwright.async_api import async_playwright
from tenacity import retry, stop_after_attempt, wait_exponential
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import functools
import hashlib

# Try to import PostgreSQL for caching
try:
    import postgres_db
    POSTGRES_AVAILABLE = postgres_db.pg_pool is not None
except ImportError:
    POSTGRES_AVAILABLE = False
    print("PostgreSQL caching not available - continuing without cache")
import builtins

# Override print to always flush
original_print = print
builtins.print = functools.partial(original_print, flush=True)

# Also ensure stdout is unbuffered
import sys
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(line_buffering=True)
sys.stdout = os.fdopen(sys.stdout.fileno(), 'w', 1)

# ────────────────────────────────────────────────────────────────────────────────
# Constants
# ────────────────────────────────────────────────────────────────────────────────
MAX_PAGES = 6
MODEL_NAME = "claude-sonnet-4-20250514"

# Cache settings
CACHE_TTL = 3600

def get_cache_key(url: str, prefix: str = "crawl") -> str:
    """Generate a cache key for a URL"""
    url_hash = hashlib.md5(url.encode()).hexdigest()
    return f"{prefix}:{url_hash}"

def cache_page(url: str, content: str, text: str) -> None:
    """Cache a crawled page if PostgreSQL is available"""
    if POSTGRES_AVAILABLE:
        try:
            cache_data = {
                'url': url,
                'content': content,
                'text': text,
                'timestamp': time.time()
            }
            postgres_db.cache_set(get_cache_key(url), cache_data, ttl=CACHE_TTL)
            print(f"  [CACHE] Saved page to cache: {url}")
        except Exception as e:
            print(f"  [CACHE] Error caching page: {e}")

def get_cached_page(url: str) -> Tuple[str, str] | None:
    """Get a page from cache if available"""
    if POSTGRES_AVAILABLE:
        try:
            cache_data = postgres_db.cache_get(get_cache_key(url))
            if cache_data:
                print(f"  [CACHE] Found page in cache: {url}")
                return cache_data.get('content', ''), cache_data.get('text', '')
        except Exception as e:
            print(f"  [CACHE] Error reading cache: {e}")
    return None

# Updated system prompt to match the PDF structure
SYSTEM_PROMPT = """
You are creating a 5-slide PowerPoint deck based on website content. You must return ONLY valid JSON.

SLIDE STRUCTURE:
---------------
Slide 0 - Title Slide:
- title: MUST be "[Company Name] and Westbridge Capital" where [Company Name] is the name of the company from the website
- subtitle: "Confidential — [Current Month Year]" (e.g., "Confidential — July 2025")

Slide 1 - Problem Overview (4 key stats):
- title: Main headline (e.g., "About Time Someone Built Intelligence into Grid Resilience")
- stat1_header: First stat header
- stat1_text: First stat description
- stat2_header: Second stat header  
- stat2_text: Second stat description
- stat3_header: Third stat header
- stat3_text: Third stat description
- stat4_header: Fourth stat header
- stat4_text: Fourth stat description
- bottom_text: Summary paragraph

Slide 2 - Solution Comparison:
- title: Slide title (e.g., "Traditional Planning vs Rhizome's Intelligence Layer")
- left_column_title: Left column header
- left_column_bullets: Array of lengthy bullet points
- right_column_title: Right column header
- right_column_bullets: Array of lengthy bullet points

Slide 3 - Market Structure (4-Box Layout):
- title: Slide title (e.g., "Some Key Observations about Market Structure")
- observations: Array of exactly 4 objects, each with:
  - heading: Short category title (e.g., "Memory Bottleneck", "Energy Consumption")
  - content: Detailed explanation (2-3 sentences about this specific market observation)

Slide 4 - Competitive Landscape Table:
- title: Slide title (e.g., "Competitive Landscape - Infrastructure Automation Market")
- table_data: Array of objects, each with:
  - vertical: Category name (e.g., "IaC Tools", "Cloud Platforms")
  - legacy: Legacy/traditional solution
  - modern: Modern/innovative solution

JSON FORMAT:
Return a JSON object with a "slides" array containing 5 objects, each with:
- "slide_number": 0-4
- Additional fields as specified above

Example:
{
  "slides": [
    {
      "slide_number": 0,
      "title": "StackGen and Westbridge Capital",
      "subtitle": "Confidential — July 2025"
    },
    {
      "slide_number": 1,
      "title": "About Time Someone Built Intelligence into Infra Provisioning",
      "stat1_header": "97% Users Report IaC Difficulties",
      "stat1_text": "56% had trouble enforcing consistent configurations...",
      "stat2_header": "56% Struggle with Consistency",
      "stat2_text": "97% of IaC users report difficulties...",
      "stat3_header": "51% of Developers Waste 20% of Time",
      "stat3_text": "51% of developers dedicate more than 20% of their time...",
      "stat4_header": "93% Demand Innovation",
      "stat4_text": "93% of respondents agree that innovation is needed...",
      "bottom_text": "StackGen's generative infra platform deterministically generates..."
    },
    {
      "slide_number": 2,
      "right_column_title": "StackGen's AI-Native Approach",
      "left_column_title": "Traditional Infrastructure as Code",
      "left_column_bullets": ["Manual template writing", "Version control complexity", "State management issues"],
      "right_column_title": "StackGen's Generative Infrastructure",
      "right_column_bullets": ["AI-powered generation", "Automatic compliance", "Intelligent optimization"]
    },
    {
      "slide_number": 3,
      "title": "Some Key Observations about Market Structure",
      "observations": [
        {
          "heading": "Memory Bottleneck",
          "content": "Over the last few years, compute power has scaled significantly faster than Memory and Bandwidth, leading to latency issues. As parameter size increases in ML models, the delay in data transfer becomes significantly longer."
        },
        {
          "heading": "Energy Consumption",
          "content": "Megatron-LM language model trained by Nvidia over nine days consumed 2.6x the annual power consumption of average U.S. house. Issues with cooling AI chips are leading to server meltdowns at some CSPs."
        },
        {
          "heading": "Profitability",
          "content": "AI compute demand scales linearly with developer headcount. At scale, GenAI companies could still face meaningful compute costs and fail to reach SaaS profitability levels."
        },
        {
          "heading": "Training vs Inference",
          "content": "One way of addressing profitability issues is to separate inference from more intensive training workload and energy-hungry chips, as it requires significantly less compute."
        }
      ]
    },
    {
      "slide_number": 4,
      "title": "Competitive Landscape - Infrastructure Automation Market",
      "table_data": [
        {
          "vertical": "IaC Tools",
          "legacy": "Terraform, CloudFormation, Pulumi",
          "modern": "StackGen, Resourcely, Firefly"
        },
        {
          "vertical": "Cloud Platforms",
          "legacy": "Manual console configuration",
          "modern": "AI-assisted provisioning"
        },
        {
          "vertical": "Security & Compliance",
          "legacy": "Post-deployment scanning",
          "modern": "Built-in policy enforcement"
        },
        {
          "vertical": "Cost Optimization",
          "legacy": "Manual review and tagging",
          "modern": "Predictive optimization"
        },
        {
          "vertical": "Multi-cloud Management",
          "legacy": "Separate tools per cloud",
          "modern": "Unified abstraction layer"
        }
      ]
    }
  ]
}

IMPORTANT: 
- The title on slide 0 MUST follow the format "[Company Name] and Westbridge Capital"
- NEVER MENTION A COMPANY'S FUNDRAISE, whatever content you add is supposed to be actual deep non-obvious insights into the company's product. 
- DO NOT INCLUDE ANY CITE TAGS.
- Return ONLY valid JSON. No markdown, no explanation, no backticks.
"""

# ────────────────────────────────────────────────────────────────────────────────
# Helper dataclasses
# ────────────────────────────────────────────────────────────────────────────────
@dataclass
class Page:
    url: str
    text: str

@dataclass
class CrawlResult:
    pages: List[Page] = field(default_factory=list)
    skipped: List[str] = field(default_factory=list)

# ────────────────────────────────────────────────────────────────────────────────
# Browserbase session helpers
# ────────────────────────────────────────────────────────────────────────────────
async def create_browserbase_context(playwright):
    """Spin up a remote Chromium session on Browserbase and return context."""
    import browserbase  # type: ignore

    bb = browserbase.Browserbase(api_key=os.environ["BROWSERBASE_API_KEY"])
    session = bb.sessions.create(project_id=os.environ["BROWSERBASE_PROJECT_ID"])
    ws_endpoint = session.connect_url
    browser = await playwright.chromium.connect_over_cdp(ws_endpoint)
    context = await browser.new_context()
    return context, session

# ────────────────────────────────────────────────────────────────────────────────
# Crawling & extraction
# ────────────────────────────────────────────────────────────────────────────────
async def fetch_page(context, url: str, stagehand: bool = False) -> str:
    """Render a page and return readable text."""
    print(f"DEBUG: fetch_page called for {url}")
    
    page = await context.new_page()
    
    # Try different wait strategies in order of preference
    wait_strategies = [
        ("domcontentloaded", 30000),  # Wait for DOM, 30s timeout
        ("load", 45000),              # Wait for load event, 45s timeout  
        (None, 60000)                 # No wait condition, 60s timeout
    ]
    
    html = None
    for wait_until, timeout in wait_strategies:
        try:
            if wait_until:
                await page.goto(url, wait_until=wait_until, timeout=timeout)
                print(f"DEBUG: Successfully navigated to {url} with {wait_until}")
            else:
                await page.goto(url, timeout=timeout)
                print(f"DEBUG: Successfully navigated to {url} with no wait condition")
            
            html = await page.content()
            break  # Success, exit the retry loop
            
        except Exception as e:
            print(f"DEBUG: Failed to navigate to {url} with {wait_until or 'no wait'}: {e}")
            if wait_until is None:  # Last attempt failed
                await page.close()
                return ""
            continue  # Try next strategy
    
    if not html:
        print(f"DEBUG: All navigation strategies failed for {url}")
        await page.close()
        return ""

    # Optional AI‑driven interactions via Stagehand
    if stagehand:
        try:
            import stagehand  # type: ignore
            sh = stagehand.Stagehand(page)
            # Example: scroll to bottom to trigger lazy‑load
            await sh.page.act("Scroll to the bottom of the page to load all content.")
            # Get updated content after Stagehand actions
            html = await page.content()
        except Exception as e:
            print(f"DEBUG: Stagehand failed: {e}")

    print(f"DEBUG: Got HTML content, length: {len(html)}")
    
    # Try multiple extraction strategies
    cleaned = None
    
    # Strategy 1: Use Trafilatura with downloaded content
    try:
        downloaded = trafilatura.fetch_url(url)
        print(f"DEBUG: Downloaded content length: {len(downloaded) if downloaded else 0}")
        if downloaded:
            cleaned = trafilatura.extract(downloaded, include_comments=False)
            print(f"DEBUG: Trafilatura extracted from download: {len(cleaned) if cleaned else 0} chars")
    except Exception as e:
        print(f"DEBUG: Trafilatura download failed: {e}")
    
    # Strategy 2: Use Trafilatura with page HTML
    if not cleaned:
        try:
            cleaned = trafilatura.extract(html, include_comments=False)
            print(f"DEBUG: Trafilatura extracted from HTML: {len(cleaned) if cleaned else 0} chars")
        except Exception as e:
            print(f"DEBUG: Trafilatura HTML extraction failed: {e}")
    
    # Strategy 3: Simple text extraction fallback
    if not cleaned:
        try:
            # Try BeautifulSoup if available
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html, 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            cleaned = soup.get_text()
            # Clean up whitespace
            lines = (line.strip() for line in cleaned.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            cleaned = ' '.join(chunk for chunk in chunks if chunk)
            print(f"DEBUG: BeautifulSoup fallback extracted: {len(cleaned)} chars")
        except ImportError:
            # BeautifulSoup not available, use simple regex fallback
            print("DEBUG: BeautifulSoup not available, using regex fallback")
            # Remove script and style tags
            import re
            text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
            # Remove all HTML tags
            text = re.sub(r'<[^>]+>', '', text)
            # Clean up whitespace
            text = re.sub(r'\s+', ' ', text)
            cleaned = text.strip()
            print(f"DEBUG: Regex fallback extracted: {len(cleaned)} chars")
        except Exception as e:
            print(f"DEBUG: Fallback extraction failed: {e}")
            cleaned = ""
    
    print(f"DEBUG: Final cleaned text length: {len(cleaned) if cleaned else 0}")
    print(f"DEBUG: First 200 chars of cleaned text: {cleaned[:200] if cleaned else 'None'}")
    
    await page.close()
    return cleaned or ""

async def fetch_page_with_html(context, url: str, stagehand: bool = False) -> Tuple[str, str]:
    """Render a page and return both readable text and HTML content."""
    print(f"DEBUG: fetch_page_with_html called for {url}")
    
    # Check cache first
    cached = get_cached_page(url)
    if cached:
        return cached
    
    page = await context.new_page()
    
    # Try different wait strategies in order of preference
    wait_strategies = [
        ("domcontentloaded", 30000),  # Wait for DOM, 30s timeout
        ("load", 45000),              # Wait for load event, 45s timeout  
        (None, 60000)                 # No wait condition, 60s timeout
    ]
    
    html = None
    for wait_until, timeout in wait_strategies:
        try:
            if wait_until:
                await page.goto(url, wait_until=wait_until, timeout=timeout)
                print(f"DEBUG: Successfully navigated to {url} with {wait_until}")
            else:
                await page.goto(url, timeout=timeout)
                print(f"DEBUG: Successfully navigated to {url} with no wait condition")
            
            html = await page.content()
            break  # Success, exit the retry loop
            
        except Exception as e:
            print(f"DEBUG: Failed to navigate to {url} with {wait_until or 'no wait'}: {e}")
            if wait_until is None:  # Last attempt failed
                await page.close()
                return "", ""
            continue  # Try next strategy
    
    if not html:
        print(f"DEBUG: All navigation strategies failed for {url}")
        await page.close()
        return "", ""

    # Optional AI‑driven interactions via Stagehand
    if stagehand:
        try:
            import stagehand  # type: ignore
            sh = stagehand.Stagehand(page)
            # Example: scroll to bottom to trigger lazy‑load
            await sh.page.act("Scroll to the bottom of the page to load all content.")
            # Get updated content after Stagehand actions
            html = await page.content()
        except Exception as e:
            print(f"DEBUG: Stagehand failed: {e}")

    print(f"DEBUG: Got HTML content, length: {len(html)}")
    
    # Try multiple extraction strategies
    cleaned = None
    
    # Strategy 1: Use Trafilatura with downloaded content
    try:
        downloaded = trafilatura.fetch_url(url)
        print(f"DEBUG: Downloaded content length: {len(downloaded) if downloaded else 0}")
        if downloaded:
            cleaned = trafilatura.extract(downloaded, include_comments=False)
            print(f"DEBUG: Trafilatura extracted from download: {len(cleaned) if cleaned else 0} chars")
    except Exception as e:
        print(f"DEBUG: Trafilatura download failed: {e}")
    
    # Strategy 2: Use Trafilatura with page HTML
    if not cleaned:
        try:
            cleaned = trafilatura.extract(html, include_comments=False)
            print(f"DEBUG: Trafilatura extracted from HTML: {len(cleaned) if cleaned else 0} chars")
        except Exception as e:
            print(f"DEBUG: Trafilatura HTML extraction failed: {e}")
    
    # Strategy 3: Simple text extraction fallback
    if not cleaned:
        try:
            # Try BeautifulSoup if available
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html, 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            cleaned = soup.get_text()
            # Clean up whitespace
            lines = (line.strip() for line in cleaned.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            cleaned = ' '.join(chunk for chunk in chunks if chunk)
            print(f"DEBUG: BeautifulSoup fallback extracted: {len(cleaned)} chars")
        except ImportError:
            # BeautifulSoup not available, use simple regex fallback
            print("DEBUG: BeautifulSoup not available, using regex fallback")
            # Remove script and style tags
            import re
            text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
            # Remove all HTML tags
            text = re.sub(r'<[^>]+>', '', text)
            # Clean up whitespace
            text = re.sub(r'\s+', ' ', text)
            cleaned = text.strip()
            print(f"DEBUG: Regex fallback extracted: {len(cleaned)} chars")
        except Exception as e:
            print(f"DEBUG: Fallback extraction failed: {e}")
            cleaned = ""
    
    print(f"DEBUG: Final cleaned text length: {len(cleaned) if cleaned else 0}")
    print(f"DEBUG: First 200 chars of cleaned text: {cleaned[:200] if cleaned else 'None'}")
    
    await page.close()
    
    # Cache the result
    text = cleaned or ""
    cache_page(url, html, text)
    
    return text, html

def discover_links(html: str, current_url: str, start_url: str) -> List[str]:
    """Extract and prioritize links from HTML content."""
    import re
    from urllib.parse import urljoin, urlparse
    
    print(f"DEBUG: discover_links called for {current_url}")
    
    # Priority keywords for important pages
    priority_keywords = [
        'blog', 'about', 'product', 'features', 'how-it-works', 'solution', 
        'case-study', 'documentation', 'docs', 'guide', 'tutorial', 'learn',
        'pricing', 'plans', 'enterprise', 'team', 'company', 'news', 'press',
        'use-case', 'workflow', 'integration', 'api', 'platform', 'overview'
    ]
    
    # Extract all links
    link_pattern = r'<a[^>]+href=["\'](.*?)["\'][^>]*>(.*?)</a>'
    matches = re.findall(link_pattern, html, re.IGNORECASE | re.DOTALL)
    
    print(f"DEBUG: Found {len(matches)} <a> tags with href attributes")
    
    found_links = []
    priority_links = []
    regular_links = []
    
    start_domain = urlparse(start_url).netloc.lower()
    print(f"DEBUG: Start domain: {start_domain}")
    
    for i, (href, link_text) in enumerate(matches):
        if i < 5:  # Debug first 5 links
            print(f"DEBUG: Processing link {i+1}: href='{href}', text='{link_text[:50]}...'")
        
        # Convert to absolute URL
        abs_url = urljoin(current_url, href)
        parsed = urlparse(abs_url)
        
        # Skip if not same domain
        if parsed.netloc.lower() != start_domain:
            if i < 5:
                print(f"DEBUG: Skipping external link: {abs_url}")
            continue
            
        # Skip fragments, mailto, tel, javascript
        if (href.startswith('#') or href.startswith('mailto:') or 
            href.startswith('tel:') or href.startswith('javascript:')):
            if i < 5:
                print(f"DEBUG: Skipping special link: {href}")
            continue
            
        # Skip common file types we don't want
        skip_extensions = ['.pdf', '.jpg', '.jpeg', '.png', '.gif', '.svg', '.css', '.js', '.xml', '.zip']
        if any(abs_url.lower().endswith(ext) for ext in skip_extensions):
            if i < 5:
                print(f"DEBUG: Skipping file link: {abs_url}")
            continue
        
        # Clean URL (remove query params and fragments for deduplication)
        clean_url = f"{parsed.scheme}://{parsed.netloc}{parsed.path}"
        if clean_url.endswith('/'):
            clean_url = clean_url[:-1]
            
        if clean_url and clean_url != current_url.rstrip('/'):
            # Check if this is a priority link
            url_text = (abs_url + ' ' + link_text).lower()
            is_priority = any(keyword in url_text for keyword in priority_keywords)
            
            if is_priority:
                priority_links.append(clean_url)
                if i < 5:
                    print(f"DEBUG: Added priority link: {clean_url}")
            else:
                regular_links.append(clean_url)
                if i < 5:
                    print(f"DEBUG: Added regular link: {clean_url}")
    
    print(f"DEBUG: Found {len(priority_links)} priority links, {len(regular_links)} regular links")
    
    # Remove duplicates while preserving order
    def dedupe(links):
        seen = set()
        result = []
        for link in links:
            if link not in seen:
                seen.add(link)
                result.append(link)
        return result
    
    priority_links = dedupe(priority_links)
    regular_links = dedupe(regular_links)
    
    # Return priority links first, then regular links (limit total)
    all_links = priority_links + regular_links
    print(f"DEBUG: Returning {len(all_links)} total unique links")
    return all_links[:50]  # Limit to prevent too many links

async def crawl_site(start_url: str, max_pages: int, stagehand: bool) -> CrawlResult:
    seen: Set[str] = set()
    queue: List[str] = [start_url]
    pages: List[Page] = []
    skipped: List[str] = []

    print(f"DEBUG: crawl_site starting with {start_url}, max_pages={max_pages}")

    async with async_playwright() as pw:
        context, session = await create_browserbase_context(pw)
        try:
            while queue and len(pages) < max_pages:
                url = queue.pop(0)
                print(f"DEBUG: Processing URL: {url}")
                
                if url in seen:
                    print(f"DEBUG: URL already seen, skipping: {url}")
                    continue
                seen.add(url)
                
                try:
                    text, html = await fetch_page_with_html(context, url, stagehand)
                    print(f"DEBUG: fetch_page returned text length: {len(text)}")
                    
                    if text:
                        pages.append(Page(url, text))
                        print(f"DEBUG: Added page {url} to pages list")
                    else:
                        print(f"DEBUG: No text extracted from {url}, not adding to pages")
                        skipped.append(f"{url} – No text content extracted")
                        
                    # Discover links from HTML content
                    if html:
                        print(f"DEBUG: Starting link discovery for {url}, HTML length: {len(html)}")
                        discovered_links = discover_links(html, url, start_url)
                        print(f"DEBUG: discover_links returned {len(discovered_links)} links: {discovered_links[:5]}...")
                        
                        added_count = 0
                        for link in discovered_links:
                            if link not in seen and link not in queue:
                                queue.append(link)
                                added_count += 1
                                print(f"DEBUG: Added link to queue: {link}")
                        
                        print(f"DEBUG: Discovered {len(discovered_links)} total links, added {added_count} new ones from {url}")
                    else:
                        print(f"DEBUG: No HTML content for link discovery from {url}")
                    
                except Exception as e:
                    print(f"DEBUG: Exception processing {url}: {e}")
                    skipped.append(f"{url} – {e}")
        finally:
            await context.close()
            # optional: bb.sessions.delete(session.id)
    
    print(f"DEBUG: crawl_site finished with {len(pages)} pages and {len(skipped)} skipped")
    return CrawlResult(pages, skipped)

# ────────────────────────────────────────────────────────────────────────────────
# Claude interaction helpers
# ────────────────────────────────────────────────────────────────────────────────
enc = tiktoken.get_encoding("cl100k_base")

def num_tokens(text: str) -> int:
    return len(enc.encode(text))

@retry(wait=wait_exponential(multiplier=2, min=4, max=30), stop=stop_after_attempt(6))
async def call_claude(client, content: str, web_search_uses: int = 10) -> dict:
    """Call Claude with the updated prompt structure."""
    message_content = [{"type": "text", "text": content}]
    
    # Add reference PDF for format guidance
    reference_pdf = "Rhizome Data and WB Capital.pdf"
    if os.path.exists(reference_pdf):
        import base64
        with open(reference_pdf, "rb") as pdf_file:
            pdf_data = base64.b64encode(pdf_file.read()).decode()
        
        message_content.append({
            "type": "document",
            "source": {
                "type": "base64",
                "media_type": "application/pdf",
                "data": pdf_data
            }
        })
    
    # Debug: Write full prompt to temp.txt
    with open("temp.txt", "w", encoding="utf-8") as f:
        f.write("="*80 + "\n")
        f.write("FULL PROMPT TO CLAUDE\n")
        f.write("="*80 + "\n\n")
        f.write("SYSTEM PROMPT:\n")
        f.write("-"*40 + "\n")
        f.write(SYSTEM_PROMPT + "\n\n")
        f.write("USER CONTENT:\n")
        f.write("-"*40 + "\n")
        f.write(content + "\n\n")
        f.write("REFERENCE PDF ATTACHED: " + ("YES - " + reference_pdf if os.path.exists(reference_pdf) else "NO") + "\n")
        f.write("MODEL: " + MODEL_NAME + "\n")
        f.write("MAX_TOKENS: 4096\n")
        f.write("TEMPERATURE: 0.2\n")
        f.write(f"WEB_SEARCH: Max {web_search_uses} uses\n\n")
    
    resp = await client.messages.create(
        model=MODEL_NAME,
        system=SYSTEM_PROMPT,
        max_tokens=8192,
        thinking={"type": "enabled",
              "budget_tokens": 4096},
        stream=False,
        messages=[{"role": "user", "content": message_content}],
        tools=[{
            "type": "web_search_20250305",
            "name": "web_search",
            "max_uses": web_search_uses
        }]
    )
    
    # Debug: Write full response to temp.txt
    with open("temp.txt", "a", encoding="utf-8") as f:
        f.write("="*80 + "\n")
        f.write("FULL RESPONSE FROM CLAUDE\n")
        f.write("="*80 + "\n\n")
        f.write("RAW RESPONSE OBJECT:\n")
        f.write("-"*40 + "\n")
        f.write(str(resp) + "\n\n")
        f.write("CONTENT BLOCKS:\n")
        f.write("-"*40 + "\n")
        
        for i, content_block in enumerate(resp.content):
            f.write(f"BLOCK {i+1} - Type: {content_block.type}\n")
            f.write(str(content_block) + "\n\n")
    
    # Handle web search responses - extract the final text content with detailed logging
    response_text = ""
    web_search_details = []
    tool_use_count = 0
    tool_result_count = 0
    
    print("\n" + "="*60)
    print("[ANALYSIS] WEB SEARCH ANALYSIS")
    print("="*60)
    
    for i, content_block in enumerate(resp.content):
        block_type = content_block.type
        print(f"[BLOCK] BLOCK {i+1}: {block_type}")
        
        if block_type == "text":
            response_text += content_block.text
            print(f"   [INFO] Text content: {len(content_block.text)} chars")
            
        elif block_type == "tool_use":
            tool_use_count += 1
            tool_name = getattr(content_block, 'name', 'unknown')
            if tool_name == "web_search":
                search_input = getattr(content_block, 'input', {})
                search_term = search_input.get('search_term', 'unknown')
                print(f"   [SEARCH] WEB SEARCH #{tool_use_count} INITIATED")
                print(f"       Search term: '{search_term}'")
                web_search_details.append({
                    'search_number': tool_use_count,
                    'search_term': search_term,
                    'status': 'initiated'
                })
            else:
                print(f"   [TOOL] Other tool use: {tool_name}")
                
        elif block_type == "tool_result":
            tool_result_count += 1
            tool_use_id = getattr(content_block, 'tool_use_id', 'unknown')
            is_error = getattr(content_block, 'is_error', False)
            content = getattr(content_block, 'content', 'No content')
            
            print(f"   [RESULT] TOOL RESULT #{tool_result_count}")
            if is_error:
                print(f"       [ERROR] ERROR: {content}")
            else:
                print(f"       [SUCCESS] SUCCESS: {len(str(content))} chars of results")
                # Try to extract search results summary
                if isinstance(content, list) and len(content) > 0:
                    first_result = content[0]
                    if hasattr(first_result, 'text'):
                        snippet = first_result.text[:100] + "..." if len(first_result.text) > 100 else first_result.text
                        print(f"       [SNIPPET] First result snippet: {snippet}")
            
            # Update corresponding search detail
            if tool_result_count <= len(web_search_details):
                web_search_details[tool_result_count-1]['status'] = 'completed' if not is_error else 'error'
                web_search_details[tool_result_count-1]['result_length'] = len(str(content))
        else:
            print(f"   [UNKNOWN] Unknown block type: {block_type}")
    
    # Summary of web search usage
    print("\n" + "="*60)
    print("[SUMMARY] WEB SEARCH SUMMARY")
    print("="*60)
    print(f"[SEARCH] Total web searches initiated: {tool_use_count}")
    print(f"[RESULT] Total tool results received: {tool_result_count}")
    
    if web_search_details:
        print("\n[DETAILS] SEARCH DETAILS:")
        for detail in web_search_details:
            status_text = "[SUCCESS]" if detail['status'] == 'completed' else "[ERROR]" if detail['status'] == 'error' else "[PENDING]"
            print(f"   {status_text} Search #{detail['search_number']}: '{detail['search_term']}'")
            if 'result_length' in detail:
                print(f"       [SIZE] Result size: {detail['result_length']} chars")
    else:
        print("[WARNING] NO WEB SEARCHES DETECTED!")
        print("   This could mean:")
        print("   - Claude didn't think web search was needed")
        print("   - API key lacks web search permissions") 
        print("   - Tool configuration issue")
    
    print("="*60 + "\n")
    
    # Enhanced debug file logging
    with open("temp.txt", "a", encoding="utf-8") as f:
        f.write("="*80 + "\n")
        f.write("WEB SEARCH ANALYSIS\n")
        f.write("="*80 + "\n\n")
        f.write(f"Web searches initiated: {tool_use_count}\n")
        f.write(f"Tool results received: {tool_result_count}\n\n")
        
        if web_search_details:
            f.write("SEARCH DETAILS:\n")
            f.write("-"*40 + "\n")
            for detail in web_search_details:
                f.write(f"Search #{detail['search_number']}: {detail['search_term']}\n")
                f.write(f"Status: {detail['status']}\n")
                if 'result_length' in detail:
                    f.write(f"Result length: {detail['result_length']} chars\n")
                f.write("\n")
        else:
            f.write("NO WEB SEARCHES DETECTED\n\n")
        
        f.write("FINAL RESPONSE TEXT:\n")
        f.write("-"*40 + "\n")
        f.write(response_text + "\n\n")
    
    print(f"DEBUG: Final combined response text length: {len(response_text)}")
    print(f"DEBUG: Full web search analysis saved to temp.txt")
    
    try:
        return json.loads(response_text)
    except json.JSONDecodeError as e:
        print(f"DEBUG: JSON parsing failed: {e}")
        print(f"DEBUG: Full response: {response_text}")
        
        # Try to extract JSON from the response if it's wrapped in markdown or other text
        import re
        json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
        if json_match:
            try:
                return json.loads(json_match.group())
            except json.JSONDecodeError:
                pass
        
        # Fallback: return a basic structure
        return {
            "slides": [
                {
                    "slide_number": 0,
                    "title": "Analysis Failed",
                    "subtitle": "JSON parsing error - check logs"
                }
            ]
        }

async def build_slide_outline(pages: List[Page], reference_file: str = None, expert_info: str = None, use_double_web_search: bool = False, target_url: str = None) -> List[dict]:
    client = anthropic.AsyncAnthropic()
    
    # Combine all page content into a single text block
    combined_text = ""
    for page in pages:
        combined_text += f"\n\n=== PAGE: {page.url} ===\n{page.text}"
    
    # Handle case where web crawling failed
    if use_double_web_search and target_url:
        combined_text = f"""FALLBACK MODE: Web crawling failed for {target_url}

INSTRUCTIONS: Since direct website crawling was unsuccessful, you must use extensive web searches to gather comprehensive information about this company. You have access to 20 web searches (double the normal amount).

Please search extensively for:
1. Company name, mission, and core business
2. Products, services, and technology stack
3. Recent news, funding, and market position
4. Industry trends and competitive landscape
5. Key statistics and market data
6. Leadership team and company culture
7. Customer base and use cases

Target URL that failed to crawl: {target_url}

Use your web search tool liberally to build a complete understanding of this company and create a compelling 5-slide presentation."""
        print(f"DEBUG: Using fallback mode with double web searches ({web_search_uses} max uses) for {target_url}")
    
    # Add expert info as additional context if provided
    if expert_info:
        combined_text += f"\n\n=== ADDITIONAL CUSTOMER INSIGHTS ===\n{expert_info}\n\nNote: The above contains expert transcripts or customer information that should be used subtly as additional context. Do not over-index on this information, but use it to enhance insights where relevant."
    
    print(f"DEBUG: build_slide_outline processing {len(pages)} pages, total text length: {len(combined_text)}")
    if use_double_web_search:
        print(f"DEBUG: FALLBACK MODE - Using {web_search_uses} web searches (double normal) for {target_url}")
    if expert_info:
        print(f"DEBUG: Including expert customer information ({len(expert_info)} chars)")
    print(f"DEBUG: First 300 chars: {combined_text[:300]}...")
    
    # Single Claude call with all content - use double web search if crawling failed
    web_search_uses = 20 if use_double_web_search else 10
    outline = await call_claude(client, combined_text, web_search_uses=web_search_uses)
    print(f"DEBUG: Claude returned: {outline}")
    
    slides = outline.get("slides", [])
    print(f"DEBUG: Extracted {len(slides)} slides from Claude response")
    
    return slides[:5]

# ────────────────────────────────────────────────────────────────────────────────
# PowerPoint builder with analysis and flexible mapping
# ────────────────────────────────────────────────────────────────────────────────

def set_text_color(text_frame, color_rgb=None, is_white=False):
    """Helper to set text color - default gray (24% brightness), optional white."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if is_white:
                run.font.color.rgb = RGBColor(255, 255, 255)  # White
            elif color_rgb:
                run.font.color.rgb = color_rgb
            else:
                # Gray: 24% brightness means RGB(61, 61, 61)
                run.font.color.rgb = RGBColor(61, 61, 61)

def analyze_template(template_path: str):
    """Analyze the template to understand its structure."""
    prs = Presentation(template_path)
    
    print("\n" + "="*80)
    print("TEMPLATE ANALYSIS")
    print("="*80)
    print(f"Total slides: {len(prs.slides)}")
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n--- SLIDE {slide_idx} ---")
        print(f"Layout name: {slide.slide_layout.name}")
        
        # List all shapes with text
        text_shapes = []
        table_shapes = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                shape_info = {
                    'index': shape_idx,
                    'name': shape.name if hasattr(shape, 'name') else 'unnamed',
                    'text': shape.text[:50] + '...' if len(shape.text) > 50 else shape.text,
                    'placeholder_type': None
                }
                
                # Check if it's a placeholder
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    shape_info['placeholder_type'] = str(ph_type)
                
                text_shapes.append(shape_info)
            
            # Check for tables
            if shape.has_table:
                table = shape.table
                table_info = {
                    'index': shape_idx,
                    'rows': len(table.rows),
                    'columns': len(table.columns)
                }
                table_shapes.append(table_info)
        
        print(f"Text shapes found: {len(text_shapes)}")
        for ts in text_shapes:
            print(f"  [{ts['index']}] {ts['name']} (type: {ts['placeholder_type']})")
            print(f"      Current text: '{ts['text']}'")
        
        if table_shapes:
            print(f"Tables found: {len(table_shapes)}")
            for ts in table_shapes:
                print(f"  [{ts['index']}] Table with {ts['rows']} rows × {ts['columns']} columns")
    
    print("="*80 + "\n")
    return prs

def convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> bool:
    """Convert PPTX to PDF using LibreOffice or fallback methods."""
    
    # Method 1: Try LibreOffice (most reliable)
    if try_libreoffice_conversion(pptx_path, pdf_path):
        return True
    
    # Method 2: Try Python libraries (last resort)
    if try_python_conversion(pptx_path, pdf_path):
        return True
    
    print("All PDF conversion methods failed")
    return False

def try_libreoffice_conversion(pptx_path: str, pdf_path: str) -> bool:
    """Try converting with LibreOffice - robust against font issues."""
    try:
        import subprocess
        import os
        
        # Get the directory containing the PPTX file
        pptx_dir = os.path.dirname(os.path.abspath(pptx_path))
        
        # Use LibreOffice to convert PPTX to PDF with font fallback options
        cmd = [
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', pptx_dir,
            # Add options to handle font issues gracefully
            '--nofirststartwizard',
            '--nologo',
            '--nolockcheck',
            pptx_path
        ]
        
        print(f"Trying LibreOffice conversion: {pptx_path} -> {pdf_path}")
        print("Note: Will proceed even if fonts (like Roboto) are missing - LibreOffice will substitute")
        
        # Run with extended timeout and capture both stdout and stderr
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        # Check if PDF was created regardless of warnings about fonts
        default_pdf = pptx_path.rsplit('.', 1)[0] + '.pdf'
        
        # LibreOffice might succeed even with font warnings
        if os.path.exists(default_pdf):
            # If the PDF was created but with a different name, rename it
            if default_pdf != pdf_path:
                import shutil
                shutil.move(default_pdf, pdf_path)
            
            print(f"Successfully converted to PDF with LibreOffice: {pdf_path}")
            if result.stderr and "font" in result.stderr.lower():
                print(f"Note: Font warnings occurred but conversion succeeded: {result.stderr}")
            return True
        else:
            print(f"LibreOffice conversion failed: {result.stderr}")
            print(f"Return code: {result.returncode}")
            return False
            
    except subprocess.TimeoutExpired:
        print("LibreOffice conversion timed out (extended to 120s)")
        return False
    except FileNotFoundError:
        print("LibreOffice not found. Install with: brew install --cask libreoffice")
        return False
    except Exception as e:
        print(f"LibreOffice conversion error: {e}")
        return False


def try_python_conversion(pptx_path: str, pdf_path: str) -> bool:
    """Try converting using Python libraries (basic fallback)."""
    try:
        print(f"Trying Python conversion: {pptx_path} -> {pdf_path}")
        
        # This is a placeholder - actual implementation would require additional libraries
        # like python-pptx + reportlab or similar
        print("Python-based PDF conversion not implemented (would require additional dependencies)")
        return False
        
    except Exception as e:
        print(f"Python conversion error: {e}")
        return False

def make_deck_flexible(outline: List[dict], template_path: str, output_path: str):
    """Create deck with flexible shape mapping."""
    
    # First, analyze the template
    prs = analyze_template(template_path)
    
    # Process each slide from the outline
    for slide_data in outline:
        slide_num = slide_data.get("slide_number", 0)
        
        if slide_num >= len(prs.slides):
            print(f"WARNING: Slide {slide_num} doesn't exist in template, skipping")
            continue
        
        slide = prs.slides[slide_num]
        print(f"\n--- UPDATING SLIDE {slide_num} ---")
        
        # Different strategies for each slide type
        if slide_num == 0:  # Title slide
            update_title_slide(slide, slide_data)
        elif slide_num == 1:  # Problem overview with stats
            update_stats_slide(slide, slide_data)
        elif slide_num == 2:  # Two-column comparison
            update_comparison_slide(slide, slide_data)
        elif slide_num == 3:  # Market structure - 4-box layout
            update_market_structure_slide(slide, slide_data)
        elif slide_num == 4:  # Competitive landscape
            update_competitive_slide(slide, slide_data)
    
    prs.save(output_path)
    print(f"\nDeck saved to: {output_path}")
    
    # Convert to PDF
    pdf_path = str(output_path).rsplit('.', 1)[0] + '.pdf'
    if convert_pptx_to_pdf(str(output_path), pdf_path):
        print(f"PDF created: {pdf_path}")
        return pdf_path
    else:
        print("PDF conversion failed, continuing with PPTX only")
        return None

def update_title_slide(slide, data):
    """Update title slide - Rectangle 3 at [1] for subtitle, Title 1 at [3] for title."""
    shapes = slide.shapes
    
    # Title is at shape index 3 (Title 1) - keep existing color (blue)
    if 'title' in data and len(shapes) > 3 and shapes[3].has_text_frame:
        shapes[3].text = data['title']
        print(f"  Set title: {data['title']}")
    
    # Subtitle is at shape index 1 (Rectangle 3) - set to gray
    if 'subtitle' in data and len(shapes) > 1 and shapes[1].has_text_frame:
        shapes[1].text = data['subtitle']
        set_text_color(shapes[1].text_frame)  # Gray text
        print(f"  Set subtitle: {data['subtitle']}")
def set_text_format(text_frame, font_size=None, color_rgb=None, is_white=False, 
                    is_light_gray=False, alignment=None):
    """Helper to set text color, font size, and alignment.
    
    Args:
        text_frame: The text frame to format
        font_size: Font size in points (e.g., 14)
        color_rgb: Custom RGB color tuple
        is_white: If True, sets text to white
        is_light_gray: If True, sets text to light gray (for citations)
        alignment: PP_ALIGN value for text alignment (e.g., PP_ALIGN.CENTER)
    """
    for paragraph in text_frame.paragraphs:
        # Set alignment at paragraph level
        if alignment is not None:
            paragraph.alignment = alignment
            
        for run in paragraph.runs:
            # Set color
            if is_white:
                run.font.color.rgb = RGBColor(255, 255, 255)  # White
            elif is_light_gray:
                run.font.color.rgb = RGBColor(128, 128, 128)  # Light gray
            elif color_rgb:
                run.font.color.rgb = color_rgb
            else:
                # Default gray: 24% brightness
                run.font.color.rgb = RGBColor(61, 61, 61)
            
            # Set font size
            if font_size:
                run.font.size = Pt(font_size)
def remove_citations(text):
    """Remove citation tags like ... from text."""
    import re
    # Remove any citation tags, keeping only the content inside
    # Pattern matches  content  and similar formats
    text = re.sub(r']*>(.*?)', r'\1', text, flags=re.DOTALL)
    # Also handle generic <cite> tags
    text = re.sub(r'<cite[^>]*>(.*?)</cite>', r'\1', text, flags=re.DOTALL)
    return text

def update_title_slide(slide, data):
    """Update title slide - Rectangle 3 at [1] for subtitle, Title 1 at [3] for title."""
    shapes = slide.shapes
    
    # Title is at shape index 3 (Title 1) - keep existing color (blue)
    if 'title' in data and len(shapes) > 3 and shapes[3].has_text_frame:
        shapes[3].text = remove_citations(data['title'])
        print(f"  Set title: {data['title']}")
    
    # Subtitle is at shape index 1 (Rectangle 3) - set to gray
    if 'subtitle' in data and len(shapes) > 1 and shapes[1].has_text_frame:
        shapes[1].text = remove_citations(data['subtitle'])
        set_text_color(shapes[1].text_frame)  # Gray text
        print(f"  Set subtitle: {data['subtitle']}")


# Updated update_stats_slide function with left alignment for title and sources
def update_stats_slide(slide, data):
    """Update slide 1 with specific shape indices based on template analysis."""
    shapes = slide.shapes
    
    # Title at index 1 (Title 2) - keep existing color and font size (32pt) - LEFT ALIGNED
    if 'title' in data and len(shapes) > 1 and shapes[1].has_text_frame:
        shapes[1].text = remove_citations(data['title'])
        # Set to left alignment
        set_text_format(shapes[1].text_frame, alignment=PP_ALIGN.LEFT)
        print(f"  Set title: {data['title'][:50]}...")
    
    # Stat headers at indices 6-9 (ee4pHeader1-4) - WHITE text on blue background
    # Stat texts at indices 2-5 (Rectangle 7,8,9,4) - GRAY text
    stat_mapping = [
        (6, 'stat1_header', 3, 'stat2_text'),  # ee4pHeader1 -> Rectangle 8 (swapped content)
        (7, 'stat2_header', 2, 'stat1_text'),  # ee4pHeader2 -> Rectangle 7 (swapped content)
        (8, 'stat3_header', 4, 'stat3_text'),  # ee4pHeader3 -> Rectangle 9
        (9, 'stat4_header', 5, 'stat4_text'),  # ee4pHeader4 -> Rectangle 4
    ]
    
    for header_idx, header_key, text_idx, text_key in stat_mapping:
        # Headers - white text on blue background
        if header_key in data and len(shapes) > header_idx and shapes[header_idx].has_text_frame:
            shapes[header_idx].text = remove_citations(data[header_key])
            set_text_format(shapes[header_idx].text_frame, font_size=14, is_white=True, 
                          alignment=PP_ALIGN.CENTER)
            print(f"  Set {header_key}: {data[header_key]}")
        
        # Text descriptions - gray text
        if text_key in data and len(shapes) > text_idx and shapes[text_idx].has_text_frame:
            shapes[text_idx].text = remove_citations(data[text_key])
            set_text_format(shapes[text_idx].text_frame, font_size=14, 
                          alignment=PP_ALIGN.CENTER)  # Default gray
            print(f"  Set {text_key}: {data[text_key][:50]}...")
    
    # Bottom text at index 10 (Rectangle 12) - gray text
    if 'bottom_text' in data and len(shapes) > 10 and shapes[10].has_text_frame:
        shapes[10].text = remove_citations(data['bottom_text'])
        set_text_format(shapes[10].text_frame, font_size=14, 
                      alignment=PP_ALIGN.CENTER)  # Default gray
        print(f"  Set bottom_text: {data['bottom_text'][:50]}...")
    
    # Sources at index 11 (ee4pContent1) - light gray text - LEFT ALIGNED
    if len(shapes) > 11 and shapes[11].has_text_frame:
        # Keep existing sources text or update if provided
        if 'sources' in data:
            shapes[11].text = remove_citations(data['sources'])
        set_text_format(shapes[11].text_frame, font_size=14, is_light_gray=True,
                      alignment=PP_ALIGN.LEFT)  # Changed to LEFT alignment
        print(f"  Set sources with light gray color, font size 14, and left alignment")

# Updated update_bullets_slide function with 16pt font
def update_bullets_slide(slide, data, slide_type):
    """Update slide 3 - title at index 2, bullets at index 1."""
    shapes = slide.shapes
    
    # Title at index 2 (Title 1) - keep existing color and size
    if 'title' in data and len(shapes) > 2 and shapes[2].has_text_frame:
        shapes[2].text = remove_citations(data['title'])
        print(f"  Set title: {data['title']}")
    
    # Bullets at index 1 (TextBox 22) - gray text, 16pt font, 1.5 line spacing
    if 'bullets' in data and len(shapes) > 1 and shapes[1].has_text_frame:
        # Clean citations from each bullet
        clean_bullets = [remove_citations(bullet) for bullet in data['bullets']]
        # Format as numbered list
        bullets_text = '\n'.join(f'{i+1}. {b}' for i, b in enumerate(clean_bullets))
        shapes[1].text = bullets_text
        
        # Apply formatting with 16pt font and 1.5 line spacing
        text_frame = shapes[1].text_frame
        for paragraph in text_frame.paragraphs:
            # Set line spacing to 1.5
            paragraph.line_spacing = 1.5
            
            for run in paragraph.runs:
                # Set gray color
                run.font.color.rgb = RGBColor(61, 61, 61)
                # Set font size to 16pt (increased from 14pt)
                run.font.size = Pt(16)
        
        print(f"  Set bullets: {len(data['bullets'])} items with 16pt font and 1.5 line spacing")

def update_market_structure_slide(slide, data):
    """Update slide 3 - Market Structure with 4-box layout based on template analysis."""
    shapes = slide.shapes
    
    # Title should already be set by the template
    if 'title' in data and hasattr(slide.shapes, 'title') and slide.shapes.title:
        slide.shapes.title.text = remove_citations(data['title'])
        print(f"  Set title: {data['title']}")
    
    # Based on our template analysis, we have 4 heading rectangles and 4 content text boxes
    # Heading rectangles: Rectangle 1, Rectangle 2, Rectangle 3, Rectangle 4
    # Content boxes: TextBox 12, TextBox 11, TextBox 13, TextBox 14
    
    # Map the observations to the correct shapes
    if 'observations' in data and len(data['observations']) >= 4:
        observations = data['observations']
        
        # Mapping based on position analysis:
        # Top-left: Rectangle 1 (heading) + TextBox 12 (content) 
        # Top-right: Rectangle 2 (heading) + TextBox 11 (content)
        # Bottom-left: Rectangle 3 (heading) + TextBox 13 (content)
        # Bottom-right: Rectangle 4 (heading) + TextBox 14 (content)
        
        shape_mappings = [
            # (heading_shape_name, content_shape_name, observation_index)
            ("Rectangle 1", "TextBox 12", 0),  # Top-left
            ("Rectangle 2", "TextBox 11", 1),  # Top-right  
            ("Rectangle 3", "TextBox 13", 2),  # Bottom-left
            ("Rectangle 4", "TextBox 14", 3),  # Bottom-right
        ]
        
        for heading_name, content_name, obs_idx in shape_mappings:
            if obs_idx < len(observations):
                obs = observations[obs_idx]
                
                # Find and update heading shape
                heading_shape = None
                content_shape = None
                
                for shape in shapes:
                    if shape.name == heading_name and shape.has_text_frame:
                        heading_shape = shape
                    elif shape.name == content_name and shape.has_text_frame:
                        content_shape = shape
                
                # Update heading
                if heading_shape and 'heading' in obs:
                    heading_shape.text = remove_citations(obs['heading'])
                    # Keep the existing blue header formatting
                    print(f"  Set {heading_name} heading: {obs['heading']}")
                
                # Update content
                if content_shape and 'content' in obs:
                    content_shape.text = remove_citations(obs['content'])
                    
                    # Apply formatting to content - gray text, appropriate font size
                    text_frame = content_shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        paragraph.line_spacing = 1.2  # Slightly tighter than bullets
                        
                        for run in paragraph.runs:
                            # Set gray color for content
                            run.font.color.rgb = RGBColor(61, 61, 61)
                            # Set font size to 12pt for content boxes
                            run.font.size = Pt(12)
                    
                    print(f"  Set {content_name} content: {obs['content'][:50]}...")

def update_comparison_slide(slide, data):
    """Update slide 3 - comparison with correct shape indices based on template analysis."""
    shapes = slide.shapes
    
    # Right column title at index 1 (Title 1) - "Traversal's Intelligence Layer"
    if 'right_column_title' in data and len(shapes) > 1 and shapes[1].has_text_frame:
        shapes[1].text = remove_citations(data['right_column_title'])
        set_text_format(shapes[1].text_frame, alignment=PP_ALIGN.CENTER)
        print(f"  Set right column title: {data['right_column_title']}")
    
    # Left column title at index 5 - "Traditional AI Security Approaches"
    if 'left_column_title' in data and len(shapes) > 5 and shapes[5].has_text_frame:
        shapes[5].text = remove_citations(data['left_column_title'])
        set_text_format(shapes[5].text_frame, alignment=PP_ALIGN.CENTER)
        print(f"  Set left column title: {data['left_column_title']}")
    
    # Left column bullets - shapes 4, 7, 8 (top to bottom on left side)
    if 'left_column_bullets' in data:
        left_bullets = [remove_citations(bullet) for bullet in data['left_column_bullets']]
        bullet_shapes = [(4, 0), (7, 1), (8, 2)]
        
        for shape_idx, bullet_idx in bullet_shapes:
            if bullet_idx < len(left_bullets) and len(shapes) > shape_idx and shapes[shape_idx].has_text_frame:
                shapes[shape_idx].text = left_bullets[bullet_idx]
                set_text_format(shapes[shape_idx].text_frame, font_size=14, 
                              alignment=PP_ALIGN.CENTER)
                print(f"  Set left bullet {bullet_idx + 1}: {left_bullets[bullet_idx][:30]}...")
    
    # Right column bullets - shapes 10, 12, 13 (top to bottom on right side)
    if 'right_column_bullets' in data:
        right_bullets = [remove_citations(bullet) for bullet in data['right_column_bullets']]
        bullet_shapes = [(10, 0), (12, 1), (13, 2)]
        
        for shape_idx, bullet_idx in bullet_shapes:
            if bullet_idx < len(right_bullets) and len(shapes) > shape_idx and shapes[shape_idx].has_text_frame:
                shapes[shape_idx].text = right_bullets[bullet_idx]
                set_text_format(shapes[shape_idx].text_frame, font_size=14, 
                              alignment=PP_ALIGN.CENTER)
                print(f"  Set right bullet {bullet_idx + 1}: {right_bullets[bullet_idx][:30]}...")

def update_competitive_slide(slide, data):
    """Update slide 4 - title and table content."""
    shapes = slide.shapes
    
    # Title at index 1 (Title 2) - keep existing color
    if 'title' in data and len(shapes) > 1 and shapes[1].has_text_frame:
        shapes[1].text = remove_citations(data['title'])
        print(f"  Set title: {data['title']}")
    
    # Look for a table in the shapes
    table_found = False
    for shape in shapes:
        if shape.has_table:
            table_found = True
            table = shape.table
            print(f"  Found table with {len(table.rows)} rows and {len(table.columns)} columns")
            
            # Check if we have table_data in the response
            if 'table_data' in data:
                table_data = data['table_data']
                
                # Fill the table (skip header row, start from row 1)
                for row_idx, row_data in enumerate(table_data):
                    if row_idx + 1 < len(table.rows):  # +1 to skip header row
                        table_row = table.rows[row_idx + 1]
                        
                        # Fill each cell in the row with cleaned text and 14pt font
                        if 'vertical' in row_data and len(table_row.cells) > 0:
                            cell = table_row.cells[0]
                            cell.text = remove_citations(row_data['vertical'])
                            # Set gray text and 14pt font
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = RGBColor(61, 61, 61)
                                    run.font.size = Pt(14)
                        
                        if 'legacy' in row_data and len(table_row.cells) > 1:
                            cell = table_row.cells[1]
                            cell.text = remove_citations(row_data['legacy'])
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = RGBColor(61, 61, 61)
                                    run.font.size = Pt(14)
                        
                        if 'modern' in row_data and len(table_row.cells) > 2:
                            cell = table_row.cells[2]
                            cell.text = remove_citations(row_data['modern'])
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = RGBColor(61, 61, 61)
                                    run.font.size = Pt(14)
                        
                        print(f"    Filled row {row_idx + 1}: {row_data.get('vertical', 'N/A')}")
            
            # Alternative: if we have bullets, convert them to table rows
            elif 'bullets' in data:
                bullets = [remove_citations(bullet) for bullet in data['bullets']]
                for row_idx, bullet in enumerate(bullets[:5]):  # Max 5 rows
                    if row_idx + 1 < len(table.rows):
                        table_row = table.rows[row_idx + 1]
                        # Parse bullet for table content
                        parts = bullet.split(':')
                        if len(parts) >= 2:
                            table_row.cells[0].text = parts[0].strip()
                            table_row.cells[1].text = parts[1].strip() if len(table_row.cells) > 1 else ""
                            # Set gray text and 14pt font
                            for cell in table_row.cells:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.color.rgb = RGBColor(61, 61, 61)
                                        run.font.size = Pt(14)
                print(f"  Converted {len(bullets)} bullets to table rows")
            break
    
    if not table_found:
        print(f"  WARNING: No table found on slide 4")
        # Fallback to adding text box
        if 'bullets' in data:
            try:
                left = Inches(0.5)
                top = Inches(2)
                width = Inches(9)
                height = Inches(4)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                clean_bullets = [remove_citations(bullet) for bullet in data['bullets']]
                text_frame.text = '\n'.join(f'• {b}' for b in clean_bullets)
                set_text_color(text_frame)  # Gray text
                print(f"  Added text box with {len(data['bullets'])} bullets")
            except Exception as e:
                print(f"  Failed to add text box: {e}")
# ────────────────────────────────────────────────────────────────────────────────
# CLI
# ────────────────────────────────────────────────────────────────────────────────

def parse_args():
    p = argparse.ArgumentParser(description="Generate a 5‑slide deck from a website")
    p.add_argument("--url", required=True, help="Root URL to crawl")
    p.add_argument("--template", default="Rhizome Data and WB Capital.pptx", help="Path to .pptx template")
    p.add_argument("--output", default="output.pptx", help="Output deck path")
    p.add_argument("--max-pages", type=int, default=MAX_PAGES, help="Max pages to crawl")
    p.add_argument("--stagehand", action="store_true", help="Use Stagehand for AI navigation")
    p.add_argument("--expert-info", help="Expert transcripts or customer information to include as additional context")
    return p.parse_args()

async def main():
    args = parse_args()

    start = time.time()
    crawl_res = await crawl_site(args.url, args.max_pages, args.stagehand)
    print(f"Crawled {len(crawl_res.pages)} pages (skipped {len(crawl_res.skipped)}) in {time.time()-start:.1f}s")

    if not crawl_res.pages:
        print("WARNING: No pages crawled successfully - proceeding with double web searches for AI content generation")
        # Create empty pages list and use double web searches
        crawl_res.pages = []
        use_double_web_search = True
    else:
        use_double_web_search = False

    outline = await build_slide_outline(crawl_res.pages, expert_info=args.expert_info, use_double_web_search=use_double_web_search, target_url=args.url)
    if len(outline) < 5:
        print(f"WARNING: Only got {len(outline)} slides from Claude (expected 5)")

    pdf_path = make_deck_flexible(outline, args.template, args.output)
    print(f"\nDeck written to {args.output}")
    if pdf_path:
        print(f"PDF created at {pdf_path}")
    else:
        print("PDF conversion was not successful")

if __name__ == "__main__":
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        print("Interrupted by user", file=sys.stderr)