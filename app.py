#!/usr/bin/env python3
"""
app.py - Westbridge Internal Tooling
Outreach PowerPoint Draft Generator with Email Functionality
"""

from flask import Flask, request, jsonify, send_file, render_template, Response, redirect, url_for, session, current_app
import subprocess
import os
import uuid
import json
import threading
import queue
import time
import sqlite3
import asyncio
from pathlib import Path
from datetime import datetime, timedelta
import re
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
import anthropic
from flask_session import Session
from msal import ConfidentialClientApplication, SerializableTokenCache
from finder import finder_service
from functools import wraps
from werkzeug.middleware.proxy_fix import ProxyFix
import base64, mimetypes, requests
import csv
import io
# Import configuration
from config import config
# Import PostgreSQL operations
import postgres_db
# Import prompts
import prompts

# Authentication temporarily disabled
def auth_required(f):
    """TEMP: Auth disabled; passthrough decorator"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        return f(*args, **kwargs)
    return decorated_function

# Database setup
def init_database():
    """Initialize PostgreSQL database and migrate from SQLite if needed"""
    # Check PostgreSQL health
    health = postgres_db.health_check()
    if health and health.get('connected'):
        print(f"[SUCCESS] PostgreSQL connected - Version: {health['version']}, Records: {health['total_keys']}")
        
        # Check if we need to migrate from SQLite
        if Path('task_metadata.db').exists() and health['total_keys'] < 10:
            print("[INFO] SQLite database found, migrating to PostgreSQL...")
            if postgres_db.migrate_from_sqlite():
                print("[SUCCESS] Migration completed successfully")
                # Optionally rename SQLite file to indicate it's been migrated
                Path('task_metadata.db').rename('task_metadata.db.migrated')
        else:
            print("[INFO] PostgreSQL already initialized")
    else:
        error_msg = health.get('error') if health else 'Unknown error'
        print(f"[ERROR] PostgreSQL connection failed: {error_msg}")
        # PostgreSQL is not available - app will have limited functionality
        print("[WARNING] PostgreSQL database not available")

def save_task_metadata(task_id, metadata):
    """Save task metadata to PostgreSQL"""
    try:
        # Extract user_id safely - check if we're in a request context
        user_id = None
        try:
            # Try to get from session if we're in a request context
            user_id = session.get('user_email')
        except RuntimeError:
            # We're outside request context (e.g., in a background thread)
            pass
        
        # Fall back to user_email in metadata if session not available
        if not user_id:
            user_id = metadata.get('user_email', metadata.get('founder_email', 'unknown'))
        
        metadata['user_id'] = user_id
        
        # Save to PostgreSQL
        if postgres_db.save_task(task_id, metadata):
            print(f"[SUCCESS] Saved metadata for task {task_id}")
        else:
            print(f"[WARNING] PostgreSQL save failed - task will not be persisted")
    except Exception as e:
        print(f"[ERROR] Error saving task metadata: {e}")

def load_task_metadata(task_id):
    """Load task metadata from PostgreSQL"""
    try:
        # Try PostgreSQL first
        task_data = postgres_db.get_task(task_id)
        if task_data:
            return task_data
        
        # Task not found in PostgreSQL
        return None
    except Exception as e:
        print(f"[ERROR] Error loading task metadata: {e}")
        return None

def load_all_task_metadata():
    """Load all task metadata from PostgreSQL"""
    try:
        # Get current user from session
        user_id = session.get('user_email', 'unknown')
        
        # Get user's tasks from PostgreSQL
        tasks = postgres_db.get_user_tasks(user_id, limit=100)
        
        # Convert to dict format for compatibility
        all_metadata = {}
        if tasks:  # Safety check in case tasks is None
            for task in tasks:
                if task:  # Extra safety check
                    all_metadata[task.get('task_id', task.get('id', ''))] = task
        
        return all_metadata
    except Exception as e:
        print(f"[ERROR] Error loading all task metadata: {e}")
        return {}

def get_task_metadata(task_id):
    """Get task metadata from memory or PostgreSQL"""
    # First try in-memory cache
    if task_id in generated_content:
        return generated_content[task_id]
    
    # If not in memory, try loading from PostgreSQL
    metadata = load_task_metadata(task_id)
    if metadata:
        # Cache it in memory for future use
        generated_content[task_id] = metadata
        print(f"[SUCCESS] Loaded task {task_id} from PostgreSQL")
        return metadata
    
    return None

# Project management functions
def get_all_projects():
    """Get all projects from PostgreSQL"""
    try:
        # Get current user from session
        user_id = session.get('user_email', 'unknown')
        
        # Get user's projects from PostgreSQL
        projects = postgres_db.get_user_projects(user_id)
        return projects if projects is not None else []
    except Exception as e:
        print(f"[ERROR] Error loading projects: {e}")
        return []

def create_project(name, description="", color="#4FD1C5"):
    """Create a new project"""
    try:
        # Get current user from session
        user_id = session.get('user_email', 'unknown')
        
        # Generate unique project ID
        project_id = str(uuid.uuid4())
        
        # Save to PostgreSQL
        project_data = {
            'name': name,
            'description': description,
            'color': color,
            'user_id': user_id
        }
        
        if postgres_db.save_project(project_id, project_data):
            return project_id
        return None
    except Exception as e:
        print(f"[ERROR] Error creating project: {e}")
        return None

def delete_project(project_id):
    """Delete a project and move its tasks to General project"""
    try:
        if str(project_id) == '1':  # Cannot delete General project
            return False
            
        # Get current user from session
        user_id = session.get('user_email', 'unknown')
        
        # Delete from PostgreSQL
        return postgres_db.delete_project(project_id, user_id)
    except Exception as e:
        print(f"[ERROR] Error deleting project: {e}")
        return False

def get_project_task_count(project_id):
    """Get the number of tasks in a project"""
    try:
        # This is now handled by PostgreSQL get_user_projects which includes task_count
        projects = get_all_projects()
        for project in projects:
            if str(project.get('id', '')) == str(project_id):
                return project.get('task_count', 0)
        return 0
    except Exception as e:
        print(f"[ERROR] Error getting project task count: {e}")
        return 0

def create_app(config_name=None):
    """Application factory pattern"""
    if config_name is None:
        config_name = os.environ.get('FLASK_ENV', 'default')
    
    app = Flask(__name__)
    
    # Load configuration
    app.config.from_object(config[config_name])
    config[config_name].init_app(app)
    
    # Configure session interface based on database availability
    if os.environ.get('DATABASE_URL'):
        # Production: Use PostgreSQL sessions
        from postgres_session import PostgreSQLSessionInterface
        app.session_interface = PostgreSQLSessionInterface()
        print("[SESSION] Using PostgreSQL session storage")
    else:
        # Development: Use filesystem sessions
        # Ensure session directory exists
        session_dir = app.config.get('SESSION_FILE_DIR', '/tmp/flask-sessions')
        os.makedirs(session_dir, exist_ok=True)
        # Configure Flask-Session for filesystem storage
        Session(app)
        print("[SESSION] Using filesystem session storage")
    
    return app

# Create the Flask app
app = create_app()

app.wsgi_app = ProxyFix(app.wsgi_app, x_proto=1, x_host=1)
app.config["PREFERRED_URL_SCHEME"] = "https"


# Initialize database
init_database()

# Start session cleanup background task for PostgreSQL sessions
if os.environ.get('DATABASE_URL'):
    def cleanup_sessions():
        """Background task to clean up expired sessions"""
        while True:
            try:
                time.sleep(3600)  # Run every hour
                postgres_db.cleanup_expired()
                print("[SESSION] Cleaned up expired sessions")
            except Exception as e:
                print(f"[SESSION] Cleanup error: {e}")
    
    cleanup_thread = threading.Thread(target=cleanup_sessions, daemon=True)
    cleanup_thread.start()
    print("[SESSION] Started session cleanup background task")

# Session refresh disabled while auth is off
@app.before_request
def refresh_session():
    return

# Store progress updates and generated content
progress_queues = {}
generated_content = {}

# Note: With PostgreSQL, we'll load user-specific data on demand rather than all at startup
# This improves performance and security by only loading what each user needs
print("[INFO] Using PostgreSQL for persistent storage - data loaded on demand per user")

def extract_company_name_from_pptx(pptx_path):
    """Extract company name from the title slide of a PowerPoint file."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        
        # Get the first slide (title slide)
        if len(prs.slides) > 0:
            title_slide = prs.slides[0]
            
            # Look for text in shapes
            for shape in title_slide.shapes:
                if shape.has_text_frame and shape.text:
                    text = shape.text.strip()
                    # Look for pattern "[Company Name] and Westbridge Capital"
                    if "and Westbridge Capital" in text:
                        company_name = text.replace(" and Westbridge Capital", "").strip()
                        # Clean the company name for filename use
                        return sanitize_filename(company_name)
        
        return None
    except Exception as e:
        print(f"Error extracting company name from PPTX: {e}")
        return None

def sanitize_filename(name):
    """Clean a string to be safe for use as a filename."""
    # Remove invalid characters for filenames
    invalid_chars = '<>:"/\\|?*'
    for char in invalid_chars:
        name = name.replace(char, '')
    
    # Replace spaces with underscores and remove extra whitespace
    name = re.sub(r'\s+', '_', name.strip())
    
    # Remove leading/trailing underscores
    name = name.strip('_')
    
    return name or "Unknown_Company"




def parse_email_response(response_text):
    """Parse the SUBJECT: and EMAIL: format response"""
    print(f"[PARSE_EMAIL] Parsing response of length: {len(response_text)}")
    print(f"[PARSE_EMAIL] First 100 chars: {response_text[:100]}")
    
    lines = response_text.strip().split('\n')
    print(f"[PARSE_EMAIL] Number of lines: {len(lines)}")
    if lines:
        print(f"[PARSE_EMAIL] First line: {lines[0]}")
    
    subject = ""
    email_content = ""
    
    parsing_email = False
    subject_found = False
    
    for i, line in enumerate(lines):
        # Check for various subject formats
        line_upper = line.strip().upper()
        if line_upper.startswith('SUBJECT:') or line_upper.startswith('**SUBJECT:**'):
            # Extract subject after the colon
            subject = line.split(':', 1)[1].strip()
            subject = subject.replace('**', '')  # Remove markdown formatting
            subject_found = True
            parsing_email = True  # Start collecting email content from next line
            continue
        elif line.strip().startswith('EMAIL:') or line.strip().upper().startswith('**EMAIL:**'):
            parsing_email = True
            continue
        elif subject_found and parsing_email:
            # Skip empty lines between subject and content
            if line.strip() == "" and not email_content:
                continue
            # Collect email content
            if email_content:
                email_content += '\n' + line
            else:
                email_content = line
    
    # If we didn't find SUBJECT: format, try alternative parsing for simple format
    if not subject and not email_content:
        # Look for a line that might be the subject (usually first non-empty line)
        for i, line in enumerate(lines):
            if line.strip():
                # If this looks like a subject line (short, might have Subject: prefix)
                if 'subject:' in line.lower() or (i == 0 and len(line) < 100):
                    if 'subject:' in line.lower():
                        subject = line.split(':', 1)[1].strip().replace('**', '')
                    else:
                        subject = line.strip()
                    # Rest is email content
                    email_content = '\n'.join(lines[i+1:]).strip()
                    break
    
    # Final fallback - if still no parsing, treat first line as subject, rest as content
    if not subject and not email_content and lines:
        subject = lines[0].strip().replace('**', '')
        if len(lines) > 1:
            email_content = '\n'.join(lines[1:]).strip()
    
    print(f"[PARSE_EMAIL] Final parsed subject: '{subject}'")
    print(f"[PARSE_EMAIL] Final parsed content length: {len(email_content)}")
    
    return subject.strip(), email_content.strip()


async def retry_email_generation(email_func, *args, max_retries=2, wait_seconds=60, **kwargs):
    """
    Retry wrapper for email generation functions.
    Tries the function up to max_retries times with wait_seconds between attempts.
    Uses fallback email if all attempts fail.
    """
    import asyncio
    import time
    
    for attempt in range(max_retries + 1):  # +1 for initial attempt
        try:
            print(f"[RETRY_EMAIL] Attempt {attempt + 1} of {max_retries + 1}")
            result = await email_func(*args, **kwargs)
            print(f"[RETRY_EMAIL] Success on attempt {attempt + 1}")
            return result
        except Exception as e:
            print(f"[RETRY_EMAIL] Attempt {attempt + 1} failed: {str(e)}")
            
            if attempt < max_retries:  # Don't wait after the last attempt
                print(f"[RETRY_EMAIL] Waiting {wait_seconds} seconds before retry...")
                await asyncio.sleep(wait_seconds)
            else:
                print(f"[RETRY_EMAIL] All {max_retries + 1} attempts failed, using fallback")
                # Return fallback email
                fallback_subject = "Partnership Opportunity with Westbridge Capital"
                fallback_email = """Dear Team,

I hope this email finds you well. I'm reaching out from Westbridge Capital regarding a potential partnership opportunity with your company.

We've been following your company's progress and are impressed by your innovative approach. We believe there could be significant synergies between our growth capital expertise and your vision.

I've attached a presentation that outlines how Westbridge Capital could support your next phase of growth. The deck includes:
- An analysis of your market position
- Our perspective on growth opportunities
- How we've helped similar companies scale

Would you be available for a brief call next week to discuss how we might work together?

Best regards,
Westbridge Capital Team"""
                return fallback_subject, fallback_email


async def generate_email(company_info, company_url, sender_name="Westbridge Capital Team", user_email=None):
    """Generate a personalized email using Claude"""
    print(f"[GENERATE_EMAIL] Starting email generation")
    print(f"[GENERATE_EMAIL] URL: {company_url}")
    print(f"[GENERATE_EMAIL] Sender name: {sender_name}")
    print(f"[GENERATE_EMAIL] Company info length: {len(company_info)} chars")
    
    client = anthropic.AsyncAnthropic()
    
    prompt = prompts.get_email_generation_prompt(company_url, company_info, sender_name)
    print(f"[GENERATE_EMAIL] Prompt created, length: {len(prompt)} chars")
    
    try:
        MODEL_NAME = "claude-sonnet-4-20250514"
        print(f"[GENERATE_EMAIL] Making API call to {MODEL_NAME}")
        
        resp = await client.messages.create(
            model=MODEL_NAME,
            system=prompts.get_email_system_prompt(user_email),
            max_tokens=8192,
            thinking={"type": "enabled",
                  "budget_tokens": 4096},
            stream=False,
            messages=[{"role": "user", "content": prompt}],
            tools=[{
                "type": "web_search_20250305",
                "name": "web_search",
                "max_uses": 6
            }]
        )
        print(f"[GENERATE_EMAIL] API call successful, processing response")
        
        # Extract text content from Claude response
        response_text = ""
        for content_block in resp.content:
            if content_block.type == "text":
                response_text += content_block.text
        
        print(f"[GENERATE_EMAIL] Extracted response text length: {len(response_text)} chars")
        print(f"[GENERATE_EMAIL] Response preview: {response_text[:200]}...")
        
        if response_text.strip():
            # Parse the response to extract subject and email
            print(f"[GENERATE_EMAIL] Parsing response to extract subject and content")
            subject, email_content = parse_email_response(response_text)
            print(f"[GENERATE_EMAIL] Parsed subject: {subject[:50]}...")
            print(f"[GENERATE_EMAIL] Parsed content length: {len(email_content)} chars")
            return subject, email_content
        else:
            print(f"[GENERATE_EMAIL] ERROR: No text content found in Claude response")
            raise Exception("No text content found in Claude response")
            
    except Exception as e:
        print(f"[GENERATE_EMAIL] ERROR: Exception occurred: {type(e).__name__}")
        print(f"[GENERATE_EMAIL] ERROR: Exception message: {e}")
        import traceback
        print(f"[GENERATE_EMAIL] ERROR: Traceback: {traceback.format_exc()}")
        raise e  # Re-raise exception to be handled by retry wrapper


def submit_to_anthropic_batch_sync(batch_job_id, batch_type, rows, sector, user_email):
    """Submit batch job to Anthropic Batch API (synchronous version)"""
    try:
        print(f"[BATCH] Starting batch submission: batch_job_id={batch_job_id}, batch_type={batch_type}, rows={len(rows)}")
        print(f"[BATCH] User: {user_email}, Sector: {sector}")
        
        # Check if API key is available
        api_key = os.environ.get('ANTHROPIC_API_KEY')
        if not api_key:
            print("[BATCH] ERROR: ANTHROPIC_API_KEY not found in environment variables")
            raise Exception("ANTHROPIC_API_KEY not configured")
        else:
            print(f"[BATCH] API key found (length: {len(api_key)})")
        
        client = anthropic.Anthropic()  # Use sync client
        print("[BATCH] Anthropic client initialized successfully")
        
        # Create batch requests based on type
        requests = []
        for i, row in enumerate(rows):
            company = row['company']
            email = row['email']
            
            if batch_type == 'founder':
                prompt = prompts.get_batch_founder_prompt(company, email, sector)
            elif batch_type == 'investor':
                prompt = prompts.get_batch_investor_prompt(company, email)
            else:  # 'any' type
                prompt = prompts.get_batch_any_prompt(company, email)
            
            # Create a valid custom_id (letters, numbers, underscores, hyphens only, max 64 chars)
            import re
            safe_company = re.sub(r'[^a-zA-Z0-9_-]', '_', company)[:20]  # Sanitize and truncate
            safe_batch_id = batch_job_id.replace('-', '_')[:20]  # UUIDs have hyphens
            custom_id = f"{safe_batch_id}_{i}_{safe_company}"[:64]  # Ensure max 64 chars
            
            request = {
                "custom_id": custom_id,
                "params": {
                    "model": "claude-sonnet-4-20250514",
                    "max_tokens": 4000,
                    "messages": [
                        {
                            "role": "user",
                            "content": prompt
                        }
                    ],
                    "tools": [
                        {
                            "type": "web_search_20250305",
                            "name": "web_search",
                            "max_uses": 8
                        }
                    ]
                }
            }
            requests.append(request)
            
            # Log each request being added
            print(f"[BATCH] Added request {i}: custom_id={custom_id}, company={company}, email={email}")
            print(f"[BATCH] Original company: '{company}' -> Safe company: '{safe_company}'")
        
        # Create the batch using the correct API endpoint
        print(f"[BATCH] Creating batch with {len(requests)} requests...")
        print(f"[BATCH] Sample request structure: {json.dumps(requests[0] if requests else 'No requests', indent=2)}")
        
        # Log the full request structure for the first request
        if requests:
            print("[BATCH] Full first request details:")
            print(f"  - custom_id: {requests[0]['custom_id']}")
            print(f"  - model: {requests[0]['params']['model']}")
            print(f"  - max_tokens: {requests[0]['params']['max_tokens']}")
            print(f"  - tools: {json.dumps(requests[0]['params']['tools'], indent=4)}")
            print(f"  - message preview: {requests[0]['params']['messages'][0]['content'][:100]}...")
        
        try:
            print("[BATCH] Calling client.messages.batches.create()...")
            batch_response = client.messages.batches.create(
                requests=requests
            )
            
            print(f"[BATCH] ✅ Successfully created Anthropic batch: {batch_response.id}")
            print(f"[BATCH] Batch status: {batch_response.processing_status}")
            print(f"[BATCH] Request counts: {batch_response.request_counts}")
            return batch_response.id
            
        except anthropic.APIError as api_error:
            print(f"[BATCH] Anthropic API Error: {type(api_error).__name__}")
            print(f"[BATCH] Error message: {str(api_error)}")
            print(f"[BATCH] Error status code: {getattr(api_error, 'status_code', 'N/A')}")
            print(f"[BATCH] Error response: {getattr(api_error, 'response', 'N/A')}")
            raise
        except Exception as create_error:
            print(f"[BATCH] Unexpected error during batch creation: {type(create_error).__name__}")
            print(f"[BATCH] Error details: {str(create_error)}")
            import traceback
            print(f"[BATCH] Traceback: {traceback.format_exc()}")
            raise
        
    except Exception as e:
        print(f"[BATCH] Error in submit_to_anthropic_batch: {type(e).__name__}")
        print(f"[BATCH] Error message: {str(e)}")
        import traceback
        print(f"[BATCH] Full traceback: {traceback.format_exc()}")
        return None

def check_batch_status_sync(batch_id):
    """Check the status of an Anthropic batch job (synchronous version)"""
    try:
        client = anthropic.Anthropic()
        batch = client.messages.batches.retrieve(batch_id)
        return batch
    except Exception as e:
        print(f"[BATCH] Error checking batch status: {e}")
        return None

def start_background_batch_monitor():
    """Start background batch monitoring in a separate thread with leader election for Fly.io"""
    import threading
    import asyncio
    import time
    import os
    from datetime import datetime, timedelta
    
    def run_monitor():
        """Run the async monitor in a thread"""
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        loop.run_until_complete(monitor_batches())
    
    async def monitor_batches():
        """Main monitoring loop with leader election"""
        print("[BATCH MONITOR] Starting background batch monitoring...")
        check_interval = 60  # Check every 60 seconds
        leader_check_interval = 30  # Check leader status every 30 seconds
        
        instance_id = os.environ.get('FLY_ALLOC_ID', f'local-{os.getpid()}')
        print(f"[BATCH MONITOR] Instance ID: {instance_id}")
        
        while True:
            try:
                # Check if this instance should be the leader
                is_leader = await check_or_become_leader(instance_id)
                
                if is_leader:
                    print(f"[BATCH MONITOR] This instance ({instance_id}) is the leader, checking batch statuses...")
                    
                    # Get all pending/submitted batch jobs that need checking
                    pending_jobs = get_pending_batch_jobs()
                    
                    if pending_jobs:
                        print(f"[BATCH MONITOR] Found {len(pending_jobs)} jobs to check")
                        await check_and_update_batches(pending_jobs)
                    else:
                        print("[BATCH MONITOR] No pending batch jobs to check")
                    
                    # Wait before next check
                    await asyncio.sleep(check_interval)
                else:
                    print(f"[BATCH MONITOR] This instance ({instance_id}) is not the leader, waiting...")
                    # Check leader status more frequently if not leader
                    await asyncio.sleep(leader_check_interval)
                
            except Exception as e:
                print(f"[BATCH MONITOR] Error in monitoring loop: {e}")
                # Wait a bit longer if there's an error to avoid rapid retries
                await asyncio.sleep(check_interval * 2)
    
    # Start monitoring thread
    monitor_thread = threading.Thread(target=run_monitor, daemon=True)
    monitor_thread.start()
    print("[BATCH MONITOR] Background monitoring thread started")

async def check_or_become_leader(instance_id):
    """Simple leader election using database"""
    try:
        conn = postgres_db.get_connection()
        try:
            with conn.cursor() as cursor:
                # Try to become leader or refresh leadership
                current_time = datetime.now()
                leader_timeout = timedelta(minutes=5)  # Leader expires after 5 minutes
                
                # Clean up old leaders first
                cursor.execute("""
                    DELETE FROM batch_monitor_leader 
                    WHERE updated_at < %s
                """, (current_time - leader_timeout,))
                
                # Try to insert or update leadership
                cursor.execute("""
                    INSERT INTO batch_monitor_leader (instance_id, updated_at)
                    VALUES (%s, %s)
                    ON CONFLICT (id) DO UPDATE SET 
                        instance_id = %s,
                        updated_at = %s
                    WHERE batch_monitor_leader.updated_at < %s
                    RETURNING instance_id
                """, (instance_id, current_time, instance_id, current_time, current_time - leader_timeout))
                
                result = cursor.fetchone()
                conn.commit()
                
                # Check if we are the leader
                cursor.execute("SELECT instance_id FROM batch_monitor_leader WHERE id = 1")
                current_leader = cursor.fetchone()
                
                is_leader = current_leader and current_leader[0] == instance_id
                return is_leader
                
        finally:
            postgres_db.return_connection(conn)
            
    except Exception as e:
        print(f"[BATCH MONITOR] Error in leader election: {e}")
        # Default to being leader if there's an error (fail-safe)
        return True

def ensure_leader_table():
    """Ensure the leader election table exists"""
    try:
        conn = postgres_db.get_connection()
        try:
            with conn.cursor() as cursor:
                cursor.execute("""
                    CREATE TABLE IF NOT EXISTS batch_monitor_leader (
                        id INTEGER PRIMARY KEY DEFAULT 1,
                        instance_id VARCHAR(255) NOT NULL,
                        updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                        CONSTRAINT single_leader CHECK (id = 1)
                    )
                """)
                conn.commit()
                print("[BATCH MONITOR] Leader election table ready")
        finally:
            postgres_db.return_connection(conn)
    except Exception as e:
        print(f"[BATCH MONITOR] Error creating leader table: {e}")

def get_pending_batch_jobs():
    """Get all batch jobs that need status checking"""
    try:
        conn = postgres_db.get_connection()
        try:
            with conn.cursor(cursor_factory=postgres_db.RealDictCursor) as cursor:
                # Get jobs that are pending/submitted and not too old (within 48 hours)
                cursor.execute("""
                    SELECT batch_id, anthropic_batch_id, user_id, batch_type, total_rows,
                           status, created_at
                    FROM batch_jobs 
                    WHERE status IN ('pending', 'submitted') 
                    AND anthropic_batch_id IS NOT NULL
                    AND created_at > NOW() - INTERVAL '48 hours'
                    ORDER BY created_at ASC
                """)
                return cursor.fetchall()
        finally:
            postgres_db.return_connection(conn)
    except Exception as e:
        print(f"[BATCH MONITOR] Error getting pending jobs: {e}")
        return []

async def check_and_update_batches(pending_jobs):
    """Check and update status for multiple batch jobs"""
    for job in pending_jobs:
        try:
            print(f"[BATCH MONITOR] Checking job {job['batch_id']}")
            
            # Check status with Anthropic
            batch_status = await check_batch_status(job['anthropic_batch_id'])
            
            if batch_status:
                await process_batch_status_update(job, batch_status)
            else:
                print(f"[BATCH MONITOR] Could not retrieve status for {job['batch_id']}")
                
        except Exception as job_error:
            print(f"[BATCH MONITOR] Error checking job {job['batch_id']}: {job_error}")

async def process_batch_status_update(job, batch_status):
    """Process status update for a single batch job"""
    try:
        batch_id = job['batch_id']
        current_status = batch_status.processing_status
        
        print(f"[BATCH MONITOR] Job {batch_id} status: {current_status}")
        
        if current_status == 'completed':
            print(f"[BATCH MONITOR] Job {batch_id} completed, retrieving results...")
            
            # Retrieve and store results
            results = []
            completed_count = 0
            
            async for result in batch_status.results():
                try:
                    result_data = {
                        'custom_id': result.custom_id,
                        'result': result.result.to_dict() if hasattr(result.result, 'to_dict') else str(result.result),
                        'status': 'completed'
                    }
                    results.append(result_data)
                    completed_count += 1
                except Exception as result_error:
                    print(f"[BATCH MONITOR] Error processing result: {result_error}")
                    results.append({
                        'custom_id': result.custom_id,
                        'result': None,
                        'status': 'failed',
                        'error': str(result_error)
                    })
            
            print(f"[BATCH MONITOR] Retrieved {completed_count} results for job {batch_id}")
            
            # Update database with results
            postgres_db.update_batch_job_status(
                batch_id, 
                'completed',
                completed_rows=completed_count,
                results=json.dumps(results)
            )
            
            # Store individual results for easy access
            await store_batch_results(job, results)
            
            print(f"[BATCH MONITOR] ✅ Job {batch_id} completed successfully")
            
        elif current_status == 'failed':
            print(f"[BATCH MONITOR] Job {batch_id} failed")
            postgres_db.update_batch_job_status(
                batch_id, 
                'failed',
                error_message=f'Anthropic batch processing failed with status: {current_status}'
            )
            
        elif current_status in ['validating', 'in_progress']:
            print(f"[BATCH MONITOR] Job {batch_id} is {current_status}")
            # Update status but keep monitoring
            postgres_db.update_batch_job_status(batch_id, current_status)
            
    except Exception as e:
        print(f"[BATCH MONITOR] Error processing status update for {job['batch_id']}: {e}")

async def store_batch_results(job, results):
    """Store individual batch results as separate task records"""
    try:
        user_id = job['user_id']
        batch_type = job['batch_type']
        
        for result_data in results:
            if result_data.get('status') == 'completed' and result_data.get('result'):
                try:
                    # Parse custom_id to get company info
                    custom_id_parts = result_data['custom_id'].split('_')
                    company_name = '_'.join(custom_id_parts[2:]) if len(custom_id_parts) > 2 else 'Unknown'
                    
                    # Create individual task record
                    task_id = str(uuid.uuid4())
                    
                    task_data = {
                        'batch_id': job['batch_id'],
                        'custom_id': result_data['custom_id'],
                        'company_name': company_name.replace('_', ' '),
                        'batch_type': batch_type,
                        'result': result_data['result'],
                        'created_at': datetime.now().isoformat()
                    }
                    
                    # Save as individual task for history viewing
                    postgres_db.save_task(
                        task_id=task_id,
                        user_id=user_id,
                        project_id='1',  # Default project
                        outreach_type=f'batch_{batch_type}',
                        status='completed',
                        task_data=task_data
                    )
                    
                except Exception as task_error:
                    print(f"[BATCH MONITOR] Error saving individual task: {task_error}")
                    
    except Exception as e:
        print(f"[BATCH MONITOR] Error storing batch results: {e}")

async def research_investor(investor_name, fund_name, outreach_context, user_email):
    """Research an investor using Claude with web search capabilities"""
    client = anthropic.AsyncAnthropic()
    
    prompt = prompts.get_investor_research_prompt(investor_name, fund_name, outreach_context)
    
    try:
        MODEL_NAME = "claude-sonnet-4-20250514"
        
        # Debug: Write full prompt to temp.txt for investor research
        with open("temp.txt", "w", encoding="utf-8") as f:
            f.write("="*80 + "\n")
            f.write("INVESTOR RESEARCH PROMPT TO CLAUDE\n")
            f.write("="*80 + "\n\n")
            f.write("SYSTEM PROMPT:\n")
            f.write("-"*40 + "\n")
            f.write(prompts.INVESTOR_RESEARCH_PROMPT + "\n\n")
            f.write("USER CONTENT:\n")
            f.write("-"*40 + "\n")
            f.write(prompt + "\n\n")
            f.write("MODEL: " + MODEL_NAME + "\n")
            f.write("MAX_TOKENS: 8192\n")
            f.write("WEB_SEARCH: Max 10 uses\n\n")
        
        resp = await client.messages.create(
            model=MODEL_NAME,
            system=prompts.INVESTOR_RESEARCH_PROMPT,
            max_tokens=8192,
            thinking={"type": "enabled",
                  "budget_tokens": 4096},
            stream=False,
            messages=[{"role": "user", "content": prompt}],
            tools=[{
                "type": "web_search_20250305",
                "name": "web_search",
                "max_uses": 20
            }]
        )
        
        # Debug: Write full response to temp.txt
        with open("temp.txt", "a", encoding="utf-8") as f:
            f.write("="*80 + "\n")
            f.write("INVESTOR RESEARCH RESPONSE FROM CLAUDE\n")
            f.write("="*80 + "\n\n")
            f.write("RAW RESPONSE OBJECT:\n")
            f.write("-"*40 + "\n")
            f.write(str(resp) + "\n\n")
            f.write("CONTENT BLOCKS:\n")
            f.write("-"*40 + "\n")
            
            for i, content_block in enumerate(resp.content):
                f.write(f"BLOCK {i+1} - Type: {content_block.type}\n")
                f.write(str(content_block) + "\n\n")
        
        # Handle web search responses - extract the final text content with detailed logging
        response_text = ""
        web_search_details = []
        tool_use_count = 0
        tool_result_count = 0
        
        print("\n" + "="*60)
        print("[ANALYSIS] INVESTOR RESEARCH WEB SEARCH ANALYSIS")
        print("="*60)
        
        for i, content_block in enumerate(resp.content):
            block_type = content_block.type
            print(f"[BLOCK] BLOCK {i+1}: {block_type}")
            
            if block_type == "text":
                response_text += content_block.text
                print(f"   [INFO] Text content: {len(content_block.text)} chars")
                
            elif block_type == "tool_use":
                tool_use_count += 1
                tool_name = getattr(content_block, 'name', 'unknown')
                if tool_name == "web_search":
                    search_input = getattr(content_block, 'input', {})
                    search_term = search_input.get('search_term', 'unknown')
                    print(f"   [SEARCH] WEB SEARCH #{tool_use_count} INITIATED")
                    print(f"       Search term: '{search_term}'")
                    web_search_details.append({
                        'search_number': tool_use_count,
                        'search_term': search_term,
                        'status': 'initiated'
                    })
                else:
                    print(f"   [TOOL] Other tool use: {tool_name}")
                    
            elif block_type == "tool_result":
                tool_result_count += 1
                tool_use_id = getattr(content_block, 'tool_use_id', 'unknown')
                is_error = getattr(content_block, 'is_error', False)
                content = getattr(content_block, 'content', 'No content')
                
                print(f"   [RESULT] TOOL RESULT #{tool_result_count}")
                if is_error:
                    print(f"       [ERROR] ERROR: {content}")
                else:
                    print(f"       [SUCCESS] SUCCESS: {len(str(content))} chars of results")
                    # Try to extract search results summary
                    if isinstance(content, list) and len(content) > 0:
                        first_result = content[0]
                        if hasattr(first_result, 'text'):
                            snippet = first_result.text[:100] + "..." if len(first_result.text) > 100 else first_result.text
                            print(f"       [SNIPPET] First result snippet: {snippet}")
                
                # Update corresponding search detail
                if tool_result_count <= len(web_search_details):
                    web_search_details[tool_result_count-1]['status'] = 'completed' if not is_error else 'error'
                    web_search_details[tool_result_count-1]['result_length'] = len(str(content))
            else:
                print(f"   [UNKNOWN] Unknown block type: {block_type}")
        
        # Summary of web search usage
        print("\n" + "="*60)
        print("[SUMMARY] INVESTOR RESEARCH WEB SEARCH SUMMARY")
        print("="*60)
        print(f"[SEARCH] Total web searches initiated: {tool_use_count}")
        print(f"[RESULT] Total tool results received: {tool_result_count}")
        
        if web_search_details:
            print("\n[DETAILS] SEARCH DETAILS:")
            for detail in web_search_details:
                status_text = "[SUCCESS]" if detail['status'] == 'completed' else "[ERROR]" if detail['status'] == 'error' else "[PENDING]"
                print(f"   {status_text} Search #{detail['search_number']}: '{detail['search_term']}'")
                if 'result_length' in detail:
                    print(f"       [SIZE] Result size: {detail['result_length']} chars")
        else:
            print("[WARNING] NO WEB SEARCHES DETECTED!")
            print("   This could mean:")
            print("   - Claude didn't think web search was needed")
            print("   - API key lacks web search permissions") 
            print("   - Tool configuration issue")
        
        print("="*60 + "\n")
        
        # Enhanced debug file logging
        with open("temp.txt", "a", encoding="utf-8") as f:
            f.write("="*80 + "\n")
            f.write("INVESTOR RESEARCH WEB SEARCH ANALYSIS\n")
            f.write("="*80 + "\n\n")
            f.write(f"Web searches initiated: {tool_use_count}\n")
            f.write(f"Tool results received: {tool_result_count}\n\n")
            
            if web_search_details:
                f.write("SEARCH DETAILS:\n")
                f.write("-"*40 + "\n")
                for detail in web_search_details:
                    f.write(f"Search #{detail['search_number']}: {detail['search_term']}\n")
                    f.write(f"Status: {detail['status']}\n")
                    if 'result_length' in detail:
                        f.write(f"Result length: {detail['result_length']} chars\n")
                    f.write("\n")
            else:
                f.write("NO WEB SEARCHES DETECTED\n\n")
            
            f.write("FINAL RESPONSE TEXT:\n")
            f.write("-"*40 + "\n")
            f.write(response_text + "\n\n")
        
        print(f"DEBUG: Final investor research response text length: {len(response_text)}")
        print(f"DEBUG: Full investor research analysis saved to temp.txt")
        
        if response_text.strip():
            return response_text.strip()
        else:
            raise Exception("No text content found in Claude response")
            
    except Exception as e:
        print(f"Error researching investor: {e}")
        return f"""• {investor_name} works at {fund_name} as an investor
• Research failed - unable to gather additional insights at this time
• Please try again or search manually for information about {investor_name}
• Common areas to research: recent blog posts, investment portfolio, speaking engagements
• Check their LinkedIn, Twitter, and fund website for recent activity
• Look for recent interviews or podcast appearances
• Review their educational background and career trajectory
• Identify companies they've invested in that relate to your outreach context
• Find any specific expertise or domain knowledge they're known for
• Research their investment philosophy and what they look for in startups"""

async def generate_any_outreach_email(person_name, organization, outreach_request, user_email, sender_name):
    """Generate a personalized email for any outreach using Claude with web search"""
    client = anthropic.AsyncAnthropic()
    
    prompt = prompts.get_any_outreach_email_prompt(person_name, organization, outreach_request, sender_name)
    
    try:
        MODEL_NAME = "claude-sonnet-4-20250514"
        
        # Debug: Write full prompt to temp.txt for any outreach
        with open("temp.txt", "w", encoding="utf-8") as f:
            f.write("="*80 + "\n")
            f.write("ANY OUTREACH EMAIL PROMPT TO CLAUDE\n")
            f.write("="*80 + "\n\n")
            f.write("SYSTEM PROMPT:\n")
            f.write("-"*40 + "\n")
            f.write(prompts.get_any_outreach_system_prompt(user_email) + "\n\n")
            f.write("USER CONTENT:\n")
            f.write("-"*40 + "\n")
            f.write(prompt + "\n\n")
            f.write("MODEL: " + MODEL_NAME + "\n")
            f.write("MAX_TOKENS: 8192\n")
            f.write("WEB_SEARCH: Max 15 uses\n\n")
        
        resp = await client.messages.create(
            model=MODEL_NAME,
            system=prompts.get_any_outreach_system_prompt(user_email),
            max_tokens=8192,
            thinking={"type": "enabled",
                  "budget_tokens": 4096},
            stream=False,
            messages=[{"role": "user", "content": prompt}],
            tools=[{
                "type": "web_search_20250305",
                "name": "web_search",
                "max_uses": 15
            }]
        )
        
        # Debug: Write full response to temp.txt
        with open("temp.txt", "a", encoding="utf-8") as f:
            f.write("="*80 + "\n")
            f.write("ANY OUTREACH EMAIL RESPONSE FROM CLAUDE\n")
            f.write("="*80 + "\n\n")
            f.write("RAW RESPONSE OBJECT:\n")
            f.write("-"*40 + "\n")
            f.write(str(resp) + "\n\n")
            f.write("CONTENT BLOCKS:\n")
            f.write("-"*40 + "\n")
            
            for i, content_block in enumerate(resp.content):
                f.write(f"BLOCK {i+1} - Type: {content_block.type}\n")
                f.write(str(content_block) + "\n\n")
        
        # Handle web search responses - extract the final text content with detailed logging
        response_text = ""
        web_search_details = []
        tool_use_count = 0
        tool_result_count = 0
        
        print("\n" + "="*60)
        print("[ANALYSIS] ANY OUTREACH EMAIL WEB SEARCH ANALYSIS")
        print("="*60)
        
        for i, content_block in enumerate(resp.content):
            block_type = content_block.type
            print(f"[BLOCK] BLOCK {i+1}: {block_type}")
            
            if block_type == "text":
                response_text += content_block.text
                print(f"   [INFO] Text content: {len(content_block.text)} chars")
                
            elif block_type == "tool_use":
                tool_use_count += 1
                tool_name = getattr(content_block, 'name', 'unknown')
                if tool_name == "web_search":
                    search_input = getattr(content_block, 'input', {})
                    search_term = search_input.get('search_term', 'unknown')
                    print(f"   [SEARCH] WEB SEARCH #{tool_use_count} INITIATED")
                    print(f"       Search term: '{search_term}'")
                    web_search_details.append({
                        'search_number': tool_use_count,
                        'search_term': search_term,
                        'status': 'initiated'
                    })
                else:
                    print(f"   [TOOL] Other tool use: {tool_name}")
                    
            elif block_type == "tool_result":
                tool_result_count += 1
                tool_use_id = getattr(content_block, 'tool_use_id', 'unknown')
                is_error = getattr(content_block, 'is_error', False)
                content = getattr(content_block, 'content', 'No content')
                
                print(f"   [RESULT] TOOL RESULT #{tool_result_count}")
                if is_error:
                    print(f"       [ERROR] ERROR: {content}")
                else:
                    print(f"       [SUCCESS] SUCCESS: {len(str(content))} chars of results")
                    # Try to extract search results summary
                    if isinstance(content, list) and len(content) > 0:
                        first_result = content[0]
                        if hasattr(first_result, 'text'):
                            snippet = first_result.text[:100] + "..." if len(first_result.text) > 100 else first_result.text
                            print(f"       [SNIPPET] First result snippet: {snippet}")
                
                # Update corresponding search detail
                if tool_result_count <= len(web_search_details):
                    web_search_details[tool_result_count-1]['status'] = 'completed' if not is_error else 'error'
                    web_search_details[tool_result_count-1]['result_length'] = len(str(content))
            else:
                print(f"   [UNKNOWN] Unknown block type: {block_type}")
        
        # Summary of web search usage
        print("\n" + "="*60)
        print("[SUMMARY] ANY OUTREACH EMAIL WEB SEARCH SUMMARY")
        print("="*60)
        print(f"[SEARCH] Total web searches initiated: {tool_use_count}")
        print(f"[RESULT] Total tool results received: {tool_result_count}")
        
        if web_search_details:
            print("\n[DETAILS] SEARCH DETAILS:")
            for detail in web_search_details:
                status_text = "[SUCCESS]" if detail['status'] == 'completed' else "[ERROR]" if detail['status'] == 'error' else "[PENDING]"
                print(f"   {status_text} Search #{detail['search_number']}: '{detail['search_term']}'")
                if 'result_length' in detail:
                    print(f"       [SIZE] Result size: {detail['result_length']} chars")
        else:
            print("[WARNING] NO WEB SEARCHES DETECTED!")
            print("   This could mean:")
            print("   - Claude didn't think web search was needed")
            print("   - API key lacks web search permissions") 
            print("   - Tool configuration issue")
        
        print("="*60 + "\n")
        
        # Enhanced debug file logging
        with open("temp.txt", "a", encoding="utf-8") as f:
            f.write("="*80 + "\n")
            f.write("ANY OUTREACH EMAIL WEB SEARCH ANALYSIS\n")
            f.write("="*80 + "\n\n")
            f.write(f"Web searches initiated: {tool_use_count}\n")
            f.write(f"Tool results received: {tool_result_count}\n\n")
            
            if web_search_details:
                f.write("SEARCH DETAILS:\n")
                f.write("-"*40 + "\n")
                for detail in web_search_details:
                    f.write(f"Search #{detail['search_number']}: {detail['search_term']}\n")
                    f.write(f"Status: {detail['status']}\n")
                    if 'result_length' in detail:
                        f.write(f"Result length: {detail['result_length']} chars\n")
                    f.write("\n")
            else:
                f.write("NO WEB SEARCHES DETECTED\n\n")
            
            f.write("FINAL RESPONSE TEXT:\n")
            f.write("-"*40 + "\n")
            f.write(response_text + "\n\n")
        
        print(f"DEBUG: Final any outreach email response text length: {len(response_text)}")
        print(f"DEBUG: Full any outreach email analysis saved to temp.txt")
        
        if response_text.strip():
            # Parse the response to extract subject and email content
            return parse_email_response(response_text.strip())
        else:
            raise Exception("No text content found in Claude response")
            
    except Exception as e:
        print(f"Error generating any outreach email: {e}")
        raise e  # Re-raise exception to be handled by retry wrapper

async def generate_investor_email(investor_name, fund_name, insights, outreach_context, user_email, sender_name):
    """Generate a personalized email for investor outreach using Claude with web search"""
    client = anthropic.AsyncAnthropic()
    
    prompt = prompts.get_investor_email_prompt(investor_name, fund_name, insights, outreach_context, sender_name)
    
    try:
        MODEL_NAME = "claude-sonnet-4-20250514"
        
        # Debug: Write full prompt to temp.txt for investor email
        with open("temp.txt", "w", encoding="utf-8") as f:
            f.write("="*80 + "\n")
            f.write("INVESTOR EMAIL PROMPT TO CLAUDE\n")
            f.write("="*80 + "\n\n")
            f.write("SYSTEM PROMPT:\n")
            f.write("-"*40 + "\n")
            f.write(prompts.get_investor_email_system_prompt(user_email) + "\n\n")
            f.write("USER CONTENT:\n")
            f.write("-"*40 + "\n")
            f.write(prompt + "\n\n")
            f.write("MODEL: " + MODEL_NAME + "\n")
            f.write("MAX_TOKENS: 8192\n")
            f.write("WEB_SEARCH: Max 10 uses\n\n")
        
        resp = await client.messages.create(
            model=MODEL_NAME,
            system=prompts.get_investor_email_system_prompt(user_email),
            max_tokens=8192,
            thinking={"type": "enabled",
                  "budget_tokens": 4096},
            stream=False,
            messages=[{"role": "user", "content": prompt}],
            tools=[{
                "type": "web_search_20250305",
                "name": "web_search",
                "max_uses": 10
            }]
        )
        
        # Handle web search responses - extract the final text content with detailed logging
        response_text = ""
        web_search_details = []
        tool_use_count = 0
        tool_result_count = 0
        
        print("\n" + "="*60)
        print("[ANALYSIS] INVESTOR EMAIL WEB SEARCH ANALYSIS")
        print("="*60)
        
        for i, content_block in enumerate(resp.content):
            block_type = content_block.type
            print(f"[BLOCK] BLOCK {i+1}: {block_type}")
            
            if block_type == "text":
                response_text += content_block.text
                print(f"   [INFO] Text content: {len(content_block.text)} chars")
                
            elif block_type == "tool_use":
                tool_use_count += 1
                tool_name = getattr(content_block, 'name', 'unknown')
                if tool_name == "web_search":
                    search_input = getattr(content_block, 'input', {})
                    search_term = search_input.get('search_term', 'unknown')
                    print(f"   [SEARCH] WEB SEARCH #{tool_use_count} INITIATED")
                    print(f"       Search term: '{search_term}'")
                    web_search_details.append({
                        'search_number': tool_use_count,
                        'search_term': search_term,
                        'status': 'initiated'
                    })
                else:
                    print(f"   [TOOL] Other tool use: {tool_name}")
                    
            elif block_type == "tool_result":
                tool_result_count += 1
                tool_use_id = getattr(content_block, 'tool_use_id', 'unknown')
                is_error = getattr(content_block, 'is_error', False)
                content = getattr(content_block, 'content', 'No content')
                
                print(f"   [RESULT] TOOL RESULT #{tool_result_count}")
                if is_error:
                    print(f"       [ERROR] ERROR: {content}")
                else:
                    print(f"       [SUCCESS] SUCCESS: {len(str(content))} chars of results")
            
            else:
                print(f"   [UNKNOWN] Unknown block type: {block_type}")
        
        # Summary of web search usage
        print("\n" + "="*60)
        print("[SUMMARY] INVESTOR EMAIL WEB SEARCH SUMMARY")
        print("="*60)
        print(f"[SEARCH] Total web searches initiated: {tool_use_count}")
        print(f"[RESULT] Total tool results received: {tool_result_count}")
        print("="*60 + "\n")
        
        print(f"DEBUG: Final investor email response text length: {len(response_text)}")
        
        if response_text.strip():
            # Parse the response to extract subject and email
            subject, email_content = parse_email_response(response_text)
            return subject, email_content
        else:
            raise Exception("No text content found in Claude response")
            
    except Exception as e:
        print(f"Error generating investor email: {e}")
        raise e  # Re-raise exception to be handled by retry wrapper

CLIENT_ID = app.config.get('OUTLOOK_CLIENT_ID')
CLIENT_SECRET = app.config.get('OUTLOOK_CLIENT_SECRET')
TENANT_ID = app.config.get('OUTLOOK_TENANT_ID')
TOKEN_PATH = Path("oauth_tokens")
TOKEN_PATH.mkdir(exist_ok=True)
GRAPH_ENDPOINT = "https://graph.microsoft.com/v1.0"

def send_batch_completion_notification(batch_data: dict) -> tuple[bool, str]:
    """
    Send a batch completion notification email to the user using their own Graph account.
    
    Args:
        batch_data (dict): Batch job data including user_id, results, etc.
        
    Returns:
        tuple[bool, str]: (success, message)
    """
    try:
        user_email = batch_data['user_id']
        batch_id = batch_data['batch_id']
        batch_type = batch_data['batch_type']
        total_rows = batch_data['total_rows']
        completed_rows = batch_data.get('completed_rows', 0)
        results = batch_data.get('results', {})
        sector = batch_data.get('sector', 'General')
        completed_at = batch_data.get('completed_at')
        
        # Calculate success statistics
        successful_count = 0
        failed_count = 0
        if results and isinstance(results, dict):
            successful_count = len([r for r in results.values() if r.get('status') == 'success'])
            failed_count = len([r for r in results.values() if r.get('status') == 'error'])
        
        # Format completion time
        completion_time = "recently"
        if completed_at:
            try:
                if isinstance(completed_at, str):
                    from datetime import datetime
                    completed_at = datetime.fromisoformat(completed_at.replace('Z', '+00:00'))
                completion_time = completed_at.strftime('%B %d, %Y at %I:%M %p')
            except:
                completion_time = "recently"
        
        # Create email subject
        subject = f"Your Westbridge {batch_type.title()} batch is complete! ({successful_count}/{total_rows} successful)"
        
        # Create email body
        body = f"""
<html>
<body style="font-family: Arial, sans-serif; line-height: 1.6; color: #333;">
    <div style="max-width: 600px; margin: 0 auto; padding: 20px;">
        <h2 style="color: #2c3e50; border-bottom: 2px solid #3498db; padding-bottom: 10px;">
            🎉 Batch Generation Complete
        </h2>
        
        <div style="background-color: #f8f9fa; padding: 20px; border-radius: 8px; margin: 20px 0;">
            <h3 style="color: #2c3e50; margin-top: 0;">Summary</h3>
            <ul style="list-style: none; padding: 0;">
                <li style="margin: 8px 0;"><strong>Batch Type:</strong> {batch_type.title()} Outreach</li>
                <li style="margin: 8px 0;"><strong>Sector:</strong> {sector}</li>
                <li style="margin: 8px 0;"><strong>Total Entries:</strong> {total_rows}</li>
                <li style="margin: 8px 0;"><strong>Successfully Processed:</strong> {successful_count}</li>
                <li style="margin: 8px 0;"><strong>Failed:</strong> {failed_count}</li>
                <li style="margin: 8px 0;"><strong>Completed:</strong> {completion_time}</li>
            </ul>
        </div>
        
        <div style="background-color: {'#d4edda' if successful_count > failed_count else '#f8d7da'}; 
                    border: 1px solid {'#c3e6cb' if successful_count > failed_count else '#f5c6cb'}; 
                    color: {'#155724' if successful_count > failed_count else '#721c24'}; 
                    padding: 15px; border-radius: 8px; margin: 20px 0;">
            <strong>
                {'✅ Batch completed successfully!' if successful_count > failed_count else '⚠️ Batch completed with some issues'}
            </strong>
        </div>
        
        <div style="margin: 30px 0;">
            <p>Your batch processing has been completed. You can view the detailed results and download the generated content by visiting your dashboard.</p>
            
            <div style="text-align: center; margin: 25px 0;">
                <a href="https://your-app-url.fly.dev" 
                   style="background-color: #3498db; color: white; padding: 12px 25px; 
                          text-decoration: none; border-radius: 5px; display: inline-block;">
                    View Results
                </a>
            </div>
        </div>
        
        <hr style="border: none; border-top: 1px solid #eee; margin: 30px 0;">
        
        <p style="color: #7f8c8d; font-size: 12px;">
            This email was sent automatically by Westbridge Capital's outreach system. 
            <br>Batch ID: {batch_id}
        </p>
    </div>
</body>
</html>
        """.strip()
        
        # Send email using the user's own Graph account
        success, message = send_graph_email(
            recipient=user_email,
            subject=subject,
            body=body,
            attachment_path=None,
            attachment_name=None,
            cc_recipients=None
        )
        
        return success, message
        
    except Exception as e:
        error_msg = f"Error creating batch notification email: {str(e)}"
        print(f"[BATCH_NOTIFICATION] {error_msg}")
        return False, error_msg

def send_graph_email(recipient, subject, body, attachment_path=None, attachment_name=None, cc_recipients=None):
    """
    Send an email via Microsoft Graph using the delegated token in the session.
    
    Args:
        recipient (str): Primary recipient email address
        subject (str): Email subject
        body (str): Email body content
        attachment_path (str, optional): Path to attachment file
        attachment_name (str, optional): Custom name for attachment
        cc_recipients (list, optional): List of CC email addresses
    """
    # 1. Get (or silently refresh) an access token
    token_result = _get_token_from_cache()
    if not token_result or "access_token" not in token_result:
        return False, "No Graph access token in session – please log in again."

    access_token = token_result["access_token"]

    # 2. Build the JSON payload
    message = {
        "message": {
            "subject": subject or "Westbridge outreach",
            "body": {"contentType": "HTML", "content": body.replace("\n", "<br>")},
            "toRecipients": [{"emailAddress": {"address": recipient}}],
        },
        "saveToSentItems": "true",
    }
    
    # Add CC recipients if provided
    if cc_recipients:
        # Ensure cc_recipients is a list and filter out empty strings
        if isinstance(cc_recipients, str):
            cc_recipients = [cc_recipients]
        
        cc_list = []
        for cc_email in cc_recipients:
            if cc_email and cc_email.strip():
                cc_list.append({"emailAddress": {"address": cc_email.strip()}})
        
        if cc_list:
            message["message"]["ccRecipients"] = cc_list

    if attachment_path and os.path.exists(attachment_path):
        # Read the file and base64 encode it
        print(f"[GRAPH] Attaching file: {attachment_path}")
        file_size = os.path.getsize(attachment_path)
        print(f"[GRAPH] File size: {file_size} bytes")
        
        mime_type, _ = mimetypes.guess_type(attachment_path)
        print(f"[GRAPH] MIME type: {mime_type}")
        
        with open(attachment_path, "rb") as f:
            content_bytes = base64.b64encode(f.read()).decode()
        
        print(f"[GRAPH] Base64 encoded size: {len(content_bytes)} chars")
        
        # Use custom name if provided, otherwise use original filename
        final_filename = attachment_name if attachment_name else Path(attachment_path).name
        
        attachment = {
            "@odata.type": "#microsoft.graph.fileAttachment",
            "name": final_filename,
            "contentType": mime_type or "application/octet-stream",
            "contentBytes": content_bytes,
        }
        message["message"]["attachments"] = [attachment]
        print(f"[GRAPH] Attachment added to message: {attachment['name']}")
    else:
        print(f"[GRAPH] No attachment - path: {attachment_path}, exists: {os.path.exists(attachment_path) if attachment_path else 'None'}")

    # 3. POST to /me/sendMail
    resp = requests.post(
        f"{GRAPH_ENDPOINT}/me/sendMail",
        headers={"Authorization": f"Bearer {access_token}",
                 "Content-Type": "application/json"},
        json=message,
        timeout=15,
    )

    if resp.status_code == 202:
        attachment_status = "with attachment" if attachment_path and os.path.exists(attachment_path) else "without attachment"
        return True, f"Email sent via Microsoft Graph {attachment_status}"
    else:
        print(f"[GRAPH] Error response: {resp.text}")
        return False, f"Graph error {resp.status_code}: {resp.text}"

def run_investor_research_with_progress(investor_name, fund_name, outreach_context, task_id, user_email, project_id=1, recipient_email=None, user_name=None, cc_recipients=None):
    """Run investor research and capture progress"""
    progress_queue = progress_queues[task_id]
    
    try:
        # Send initial progress
        progress_queue.put({
            'stage': 'crawling',
            'status': 'in_progress',
            'message': f'Starting research on {investor_name}...',
            'pages_crawled': 0
        })
        
        # Stage 1: Research preparation
        progress_queue.put({
            'stage': 'analyzing',
            'status': 'in_progress',
            'message': 'AI analyzing research requirements...',
            'pages_crawled': 0
        })
        
        # Stage 2: Conducting research
        progress_queue.put({
            'stage': 'generating',
            'status': 'in_progress',
            'message': 'Conducting web research and analysis...',
            'pages_crawled': 0
        })
        
        # Use asyncio to run the async research function
        try:
            import asyncio
            insights = asyncio.run(research_investor(investor_name, fund_name, outreach_context, user_email))
        except Exception as research_error:
            print(f"Error during research_investor call: {str(research_error)}")
            print(f"Research parameters: investor={investor_name}, fund={fund_name}")
            raise RuntimeError(f"Research failed: {str(research_error)}") from research_error
        
        # Count insights
        try:
            insight_lines = [line.strip() for line in insights.split('\n') if line.strip() and line.strip().startswith('•')]
            insights_count = len(insight_lines)
        except Exception as parsing_error:
            print(f"Error parsing insights: {str(parsing_error)}")
            print(f"Raw insights content: {insights[:500]}...")  # Log first 500 chars
            raise RuntimeError("Failed to parse research insights") from parsing_error
        
        # Stage 3: Generate email based on insights
        progress_queue.put({
            'stage': 'email',
            'status': 'in_progress',
            'message': f'Generated {insights_count} insights, now generating email...',
            'pages_crawled': 0
        })
        
        # Generate email based on insights
        try:
            # Use the user_name passed as parameter
            sender_name = user_name or "Westbridge Capital Team"
            email_subject, email_content = asyncio.run(retry_email_generation(generate_investor_email, investor_name, fund_name, insights, outreach_context, user_email, sender_name))
        except Exception as email_error:
            print(f"Error during generate_investor_email call: {str(email_error)}")
            print(f"Email parameters: investor={investor_name}, fund={fund_name}")
            raise RuntimeError(f"Email generation failed: {str(email_error)}") from email_error
        
        # Store generated content
        try:
            metadata = {
                'project_id': project_id,
                'file_id': None,  # No file for investor research
                'investor_insights': insights,
                'insights_count': insights_count,
                'investor_name': investor_name,
                'fund_name': fund_name,
                'email_subject': email_subject,
                'email_content': email_content,
                'user_email': user_email,
                'recipient_email': recipient_email,
                'cc_recipients': cc_recipients,
                'pptx_path': None,
                'pdf_path': None,
                'outreach_type': 'investor'
            }
            generated_content[task_id] = metadata
            
            # Also save to database for persistence
            save_task_metadata(task_id, metadata)
            
        except Exception as storage_error:
            print(f"Error storing generated content: {str(storage_error)}")
            print(f"Task ID: {task_id}")
            print(f"Content size: {len(insights)} chars")
            raise RuntimeError("Failed to store research results") from storage_error
        
        # Success
        progress_queue.put({
            'stage': 'completed',
            'status': 'completed',
            'message': 'Research and email completed!',
            'pages_crawled': 0,
            'file_id': None,
            'investor_insights': insights,
            'insights_count': insights_count,
            'email_subject': email_subject,
            'email_content': email_content,
            'user_email': user_email,
            'recipient_email': recipient_email
        })
        
    except Exception as e:
        error_msg = f"Research failed: {str(e)}"
        print(f"Error in run_investor_research_with_progress: {error_msg}")
        print(f"Full error details: {type(e).__name__}: {str(e)}")
        if hasattr(e, '__cause__'):
            print(f"Caused by: {e.__cause__}")
        
        progress_queue.put({
            'stage': 'error',
            'status': 'error',
            'message': error_msg,
            'pages_crawled': 0
        })
    finally:
        # Send end signal
        progress_queue.put(None)

def run_any_outreach_with_progress(person_name, organization, outreach_request, task_id, user_email, project_id=1, recipient_email=None, user_name=None, cc_recipients=None):
    """Run any outreach email generation and capture progress"""
    progress_queue = progress_queues[task_id]
    
    try:
        # Send initial progress
        progress_queue.put({
            'stage': 'crawling',
            'status': 'in_progress',
            'message': f'Starting research on {person_name} at {organization}...',
            'pages_crawled': 0
        })
        
        # Stage 1: Research preparation
        progress_queue.put({
            'stage': 'analyzing',
            'status': 'in_progress',
            'message': 'AI analyzing outreach requirements...',
            'pages_crawled': 0
        })
        
        # Stage 2: Conducting research and generating email
        progress_queue.put({
            'stage': 'generating',
            'status': 'in_progress',
            'message': 'Conducting web research and generating email...',
            'pages_crawled': 0
        })
        
        # Use asyncio to run the async email generation function
        try:
            import asyncio
            # Use the user_name passed as parameter
            sender_name = user_name or "Westbridge Capital Team"
            email_subject, email_content = asyncio.run(retry_email_generation(generate_any_outreach_email, person_name, organization, outreach_request, user_email, sender_name))
        except Exception as email_error:
            print(f"Error during generate_any_outreach_email call: {str(email_error)}")
            print(f"Email parameters: person={person_name}, org={organization}")
            raise RuntimeError(f"Email generation failed: {str(email_error)}") from email_error
        
        # Final stage
        progress_queue.put({
            'stage': 'email',
            'status': 'in_progress',
            'message': 'Email generated successfully...',
            'pages_crawled': 0
        })
        
        # Store generated content
        try:
            metadata = {
                'project_id': project_id,
                'file_id': None,  # No file for any outreach
                'email_subject': email_subject,
                'email_content': email_content,
                'person_name': person_name,
                'organization': organization,
                'outreach_request': outreach_request,
                'user_email': user_email,
                'recipient_email': recipient_email,
                'cc_recipients': cc_recipients,
                'pptx_path': None,
                'pdf_path': None,
                'outreach_type': 'any'
            }
            generated_content[task_id] = metadata
            
            # Also save to database for persistence
            save_task_metadata(task_id, metadata)
            
        except Exception as storage_error:
            print(f"Error storing generated content: {str(storage_error)}")
            print(f"Task ID: {task_id}")
            print(f"Content size: {len(email_content)} chars")
            raise RuntimeError("Failed to store email results") from storage_error
        
        # Success
        progress_queue.put({
            'stage': 'completed',
            'status': 'completed',
            'message': 'Email generation completed!',
            'pages_crawled': 0,
            'file_id': None,
            'email_subject': email_subject,
            'email_content': email_content,
            'person_name': person_name,
            'organization': organization,
            'user_email': user_email,
            'recipient_email': recipient_email
        })
        
    except Exception as e:
        error_msg = f"Email generation failed: {str(e)}"
        print(f"Error in run_any_outreach_with_progress: {error_msg}")
        print(f"Full error details: {type(e).__name__}: {str(e)}")
        if hasattr(e, '__cause__'):
            print(f"Caused by: {e.__cause__}")
        
        progress_queue.put({
            'stage': 'error',
            'status': 'error',
            'message': error_msg,
            'pages_crawled': 0
        })
    finally:
        # Send end signal
        progress_queue.put(None)

def run_conversion_with_progress(url, output_path, task_id, user_email, expert_info=None, project_id=1, recipient_email=None, user_name=None, cc_recipients=None):
    """Run the conversion script and capture progress"""
    progress_queue = progress_queues[task_id]
    
    try:
        # Send initial progress
        progress_queue.put({
            'stage': 'crawling',
            'status': 'in_progress',
            'message': 'Starting web crawler...',
            'pages_crawled': 0
        })
        
        # Run the PowerPoint generation script
        cmd = [
            '/app/.venv/bin/python', '-u',
            'website_to_ppt.py',
            '--url', url,
            '--output', str(output_path)
        ]
        
        # Add template if exists
        template_path = app.config['TEMPLATE_PPTX']
        if os.path.exists(template_path):
            cmd.extend(['--template', template_path])
        
        # Add expert info if provided
        if expert_info:
            cmd.extend(['--expert-info', expert_info])
        
        env = os.environ.copy()
        env['PYTHONUNBUFFERED'] = '1'
        
        process = subprocess.Popen(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
            bufsize=1,
            universal_newlines=True,
            env=env
        )
        
        pages_crawled = 0
        current_stage = 'crawling'
        company_info = []
        
        print(f"[FOUNDER EMAIL] Starting web scraping process for {url}")
        
        # Monitor output line by line
        while True:
            line = process.stdout.readline()
            if not line:
                break
                
            line = line.strip()
            if line:
                print(f"Script output: {line}", flush=True)
                
                # Collect company information from the crawled pages
                if "DEBUG: First 200 chars of cleaned text:" in line:
                    extracted_text = line.split(":", 2)[2].strip()  # Split on 2 colons to skip "DEBUG: First 200 chars of cleaned text"
                    company_info.append(extracted_text)
                    print(f"[FOUNDER EMAIL] Extracted company text snippet {len(company_info)}: {extracted_text[:100]}...")
                
                # Parse progress from script output
                if "Processing URL:" in line:
                    pages_crawled += 1
                    progress_queue.put({
                        'stage': 'crawling',
                        'status': 'in_progress',
                        'message': f'Crawling page {pages_crawled}...',
                        'pages_crawled': pages_crawled
                    })
                
                elif "build_slide_outline" in line:
                    current_stage = 'analyzing'
                    progress_queue.put({
                        'stage': 'analyzing',
                        'status': 'in_progress',
                        'message': 'AI analyzing content...',
                        'pages_crawled': pages_crawled
                    })
                
                elif "UPDATING SLIDE" in line:
                    current_stage = 'generating'
                    slide_match = re.search(r'SLIDE (\d+)', line)
                    if slide_match:
                        slide_num = int(slide_match.group(1))
                        progress_queue.put({
                            'stage': 'generating',
                            'status': 'in_progress',
                            'message': f'Creating slide {slide_num + 1} of 5...',
                            'pages_crawled': pages_crawled
                        })
                
                elif "Crawled" in line and "pages" in line:
                    match = re.search(r'Crawled (\d+) pages', line)
                    if match:
                        pages_crawled = int(match.group(1))
        
        process.wait()
        
        print(f"[FOUNDER EMAIL] Web scraping completed with return code: {process.returncode}")
        print(f"[FOUNDER EMAIL] Total company info snippets collected: {len(company_info)}")
        if company_info:
            print(f"[FOUNDER EMAIL] First snippet preview: {company_info[0][:100]}...")
        else:
            print(f"[FOUNDER EMAIL] WARNING: No company information was extracted!")
        
        if process.returncode == 0 and output_path.exists():
            # PowerPoint generated successfully, store in database for persistence
            try:
                with open(output_path, 'rb') as f:
                    pptx_data = f.read()
                postgres_db.store_generated_file(
                    file_id=task_id, 
                    user_email=user_email, 
                    filename=f"{company_name}_and_Westbridge_Capital.pptx" if company_name else "westbridge_outreach_deck.pptx",
                    file_type='pptx',
                    file_data=pptx_data,
                    company_name=company_name
                )
                print(f"[FILE STORAGE] Stored PPTX in database for task {task_id}")
            except Exception as e:
                print(f"[FILE STORAGE] Warning: Failed to store PPTX in database: {e}")
            
            # PowerPoint generated successfully, now check for PDF and generate email
            pdf_path = Path(str(output_path).rsplit('.', 1)[0] + '.pdf')
            
            progress_queue.put({
                'stage': 'email',
                'status': 'in_progress',
                'message': 'Generating personalized email...',
                'pages_crawled': pages_crawled
            })
            
            # Extract company name from URL
            from urllib.parse import urlparse
            parsed_url = urlparse(url)
            company_name = parsed_url.netloc.replace('www.', '').split('.')[0].title()
            
            # Generate email content
            company_text = '\n'.join(company_info[:5])  # Use first 5 text snippets
            import asyncio
            sender_name = user_name or "Westbridge Capital Team"
            
            print(f"[FOUNDER EMAIL] Starting email generation for {url}")
            print(f"[FOUNDER EMAIL] Company text length: {len(company_text)} chars")
            print(f"[FOUNDER EMAIL] Sender name: {sender_name}")
            
            try:
                email_subject, email_content = asyncio.run(retry_email_generation(generate_email, company_text, url, sender_name, user_email))
                print(f"[FOUNDER EMAIL] Email generation successful")
                print(f"[FOUNDER EMAIL] Subject: {email_subject[:50]}...")
                print(f"[FOUNDER EMAIL] Content length: {len(email_content)} chars")
            except Exception as email_gen_error:
                print(f"[FOUNDER EMAIL] ERROR: Email generation failed: {email_gen_error}")
                print(f"[FOUNDER EMAIL] ERROR: Company text preview: {company_text[:200]}...")
                raise
            
            # Wait for PDF conversion to complete (with timeout)
            pdf_ready = False
            pdf_wait_attempts = 0
            max_pdf_wait = 30  # Wait up to 30 seconds for PDF
            
            while pdf_wait_attempts < max_pdf_wait and not pdf_ready:
                if pdf_path.exists():
                    # Check if PDF file is complete by trying to get its size
                    try:
                        pdf_size = pdf_path.stat().st_size
                        if pdf_size > 0:
                            # Wait a bit more to ensure file is fully written
                            time.sleep(1)
                            new_size = pdf_path.stat().st_size
                            if new_size == pdf_size:  # File size stable
                                pdf_ready = True
                                # Store PDF in database for persistence
                                try:
                                    with open(pdf_path, 'rb') as f:
                                        pdf_data = f.read()
                                    postgres_db.store_generated_file(
                                        file_id=task_id, 
                                        user_email=user_email, 
                                        filename=f"{company_name}_and_Westbridge_Capital.pdf" if company_name else "westbridge_outreach_deck.pdf",
                                        file_type='pdf',
                                        file_data=pdf_data,
                                        company_name=company_name
                                    )
                                    print(f"[FILE STORAGE] Stored PDF in database for task {task_id}")
                                except Exception as e:
                                    print(f"[FILE STORAGE] Warning: Failed to store PDF in database: {e}")
                                break
                    except:
                        pass
                
                time.sleep(1)
                pdf_wait_attempts += 1
                
                # Update progress to show PDF conversion status
                if pdf_wait_attempts % 5 == 0:  # Every 5 seconds
                    progress_queue.put({
                        'stage': 'email',
                        'status': 'in_progress',
                        'message': f'Converting to PDF... ({pdf_wait_attempts}s)',
                        'pages_crawled': pages_crawled
                    })
            
            # Store generated content
            file_id = output_path.stem
            print(f"[FOUNDER EMAIL] Storing metadata for task {task_id}")
            print(f"[FOUNDER EMAIL] File ID: {file_id}")
            print(f"[FOUNDER EMAIL] Email subject exists: {bool(email_subject)}")
            print(f"[FOUNDER EMAIL] Email content exists: {bool(email_content)}")
            
            metadata = {
                'project_id': project_id,
                'file_id': file_id,
                'email_subject': email_subject,
                'email_content': email_content,
                'company_name': company_name,
                'user_email': user_email,
                'recipient_email': recipient_email,
                'cc_recipients': cc_recipients,
                'pptx_path': str(output_path.absolute()),
                'pdf_path': str(pdf_path.absolute()) if pdf_path.exists() else None,
                'slide_title': None,  # Will be populated if available
                'outreach_type': 'founder'
            }
            generated_content[task_id] = metadata
            print(f"[FOUNDER EMAIL] Metadata stored successfully")
            
            # Also save to database for persistence
            save_task_metadata(task_id, metadata)
            
            # Success
            print(f"[FOUNDER EMAIL] Sending completion status to frontend")
            print(f"[FOUNDER EMAIL] Subject for frontend: {email_subject[:50] if email_subject else 'None'}...")
            print(f"[FOUNDER EMAIL] Content length for frontend: {len(email_content) if email_content else 0}")
            
            progress_queue.put({
                'stage': 'completed',
                'status': 'completed',
                'message': 'Done!',
                'pages_crawled': pages_crawled,
                'file_id': file_id,
                'email_subject': email_subject,
                'email_content': email_content,
                'user_email': user_email,
                'recipient_email': recipient_email,
                'has_pdf': pdf_path.exists()
            })
            print(f"[FOUNDER EMAIL] Completion status sent successfully")
        else:
            # Error
            error_msg = 'Failed to generate presentation'
            if process.returncode != 0:
                error_msg += f' (exit code: {process.returncode})'
            
            progress_queue.put({
                'stage': 'error',
                'status': 'error',
                'message': error_msg,
                'pages_crawled': pages_crawled
            })
            
    except Exception as e:
        progress_queue.put({
            'stage': 'error',
            'status': 'error',
            'message': str(e)
        })
    
    finally:
        # Send end signal
        progress_queue.put(None)

# # ─── config values ─────────────────────────────────────────────────────────────
# app.config.update(
#     AUTHORITY = f"https://login.microsoftonline.com/{app.config['OUTLOOK_TENANT_ID']}",
#     OUTLOOK_SCOPES = ["openid", "profile", "email", "offline_access", "Mail.Send"],
#     REDIRECT_PATH = "/auth/callback",          # <‑‑ can be any path you like
# )

# ─── MSAL helper functions ─────────────────────────────────────────────────────

def _build_msal_app(cache=None):
    return ConfidentialClientApplication(
        current_app.config["OUTLOOK_CLIENT_ID"],
        authority=current_app.config["AUTHORITY"],
        client_credential=current_app.config["OUTLOOK_CLIENT_SECRET"],
        token_cache=cache,
    )

def _get_token_from_cache():
    cache = SerializableTokenCache()
    if (serialized := session.get("token_cache")):
        cache.deserialize(serialized)
    cca = _build_msal_app(cache)
    accounts = cca.get_accounts()
    if accounts:
        result = cca.acquire_token_silent(current_app.config["OUTLOOK_SCOPES"], accounts[0])
        # Update session cache
        session["token_cache"] = cache.serialize()
        return result

# ─── routes ────────────────────────────────────────────────────────────────────
@app.route("/login")
def login():
    # TEMP: Auth disabled; pretend login succeeded
    session["user_authenticated"] = True
    session["user_name"] = session.get("user_name", "Developer")
    session["user_email"] = session.get("user_email", "dev@example.com")
    session.permanent = True
    return redirect(url_for("index"))

@app.route("/auth/callback") #  <-- /auth/callback
def auth_callback():
    # TEMP: Auth disabled; just redirect to index
    return redirect(url_for("index"))

@app.route("/logout")
def logout():
    # TEMP: Auth disabled; just clear local flags and go home
    session.permanent = False
    session.clear()
    return redirect(url_for("index"))

@app.route("/auth/status")
def auth_status():
    """TEMP: Always report authenticated while auth is disabled"""
    return jsonify({
        "authenticated": True,
        "user_name": session.get("user_name", "Developer"),
        "user_email": session.get("user_email", "dev@example.com")
    })


@app.route('/')
def index():
    """TEMP: Always show app while auth is disabled"""
    user_info = {
        'authenticated': True,
        'name': session.get('user_name', 'Developer'),
        'email': session.get('user_email', 'dev@example.com')
    }
    return render_template('index.html', user=user_info)

@app.route('/analytics')
@auth_required
def analytics():
    """Analytics dashboard route"""
    # Pass user info to template to avoid flash of sign-in button
    user_info = {
        'authenticated': True,
        'name': session.get('user_name', 'User'),
        'email': session.get('user_email', '')
    }
    return render_template('analytics.html', user=user_info)

@app.route('/batch-jobs')
@auth_required
def batch_jobs():
    """Batch jobs management page"""
    # Pass user info to template to avoid flash of sign-in button
    user_info = {
        'authenticated': True,
        'name': session.get('user_name', 'User'),
        'email': session.get('user_email', '')
    }
    return render_template('batch_jobs.html', user=user_info)

@app.route('/history')
def history_page():
    """History viewing page"""
    # TEMP: auth disabled; allow access
    # Pass user info to template to avoid flash of sign-in button
    user_info = {
        'authenticated': True,
        'name': session.get('user_name', 'User'),
        'email': session.get('user_email', '')
    }
    return render_template('history.html', user=user_info)

@app.route('/history/<task_id>')
def view_history_item(task_id):
    """View specific history item details"""
    try:
        # TEMP: auth disabled; allow access
        # Try to get user email from session
        user_email = session.get('user_email', 'test@example.com')
        
        # Try to get task or batch job details from database
        task_data = postgres_db.get_task_or_batch_by_id(task_id, user_email)
        
        if not task_data:
            # If no database data, use mock data for testing
            from datetime import datetime
            task_data = {
                'task_id': task_id,
                'user_id': user_email,
                'outreach_type': 'founder',
                'status': 'completed',
                'created_at': datetime(2025, 7, 28, 10, 0, 0),
                'updated_at': datetime(2025, 7, 28, 10, 30, 0),
                'task_data': {
                    'url': 'https://example.com',
                    'sector': 'infra',
                    'recipient_email': 'test@example.com',
                    'file_id': task_id,
                    'email_content': f'This is test email content for task {task_id}.',
                    'generated_slides': [f'Slide 1 for {task_id}', f'Slide 2 for {task_id}']
                }
            }
        
        return render_template('history_item.html', task=task_data)
        
    except Exception as e:
        print(f"[HISTORY] Error viewing history item: {e}")
        print(f"[DEBUG] Session user_email: {session.get('user_email')}")
        print(f"[DEBUG] Task ID requested: {task_id}")
        return render_template('history_item.html', error=f'Error loading task details: {str(e)}')

@app.route('/customization')
@auth_required
def customization_page():
    """Show customization page for email templates"""
    user_info = {
        'authenticated': True,
        'name': session.get('user_name', 'User'),
        'email': session.get('user_email', '')
    }
    return render_template('customization.html', user=user_info)

@app.route('/api/email-templates', methods=['GET'])
@auth_required
def get_email_templates():
    """Get all email templates for current user"""
    try:
        user_email = session.get('user_email')
        templates = postgres_db.get_user_email_templates(user_email)
        return jsonify(templates)
    except Exception as e:
        print(f"[CUSTOMIZATION] Error getting templates: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/email-templates', methods=['POST'])
@auth_required
def add_email_template():
    """Add new email template"""
    try:
        data = request.get_json()
        user_email = session.get('user_email')
        
        template_id = postgres_db.add_email_template(
            user_email=user_email,
            template_type=data['type'],
            company_name=data['company_name'],
            email_content=data['email_content']
        )
        
        return jsonify({'id': template_id, 'message': 'Template added successfully'})
    except Exception as e:
        print(f"[CUSTOMIZATION] Error adding template: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/email-templates/<template_id>', methods=['PUT'])
@auth_required
def update_email_template(template_id):
    """Update existing email template"""
    try:
        data = request.get_json()
        user_email = session.get('user_email')
        
        success = postgres_db.update_email_template(
            template_id=template_id,
            user_email=user_email,
            template_type=data['type'],
            company_name=data['company_name'],
            email_content=data['email_content']
        )
        
        if success:
            return jsonify({'message': 'Template updated successfully'})
        else:
            return jsonify({'error': 'Template not found or unauthorized'}), 404
    except Exception as e:
        print(f"[CUSTOMIZATION] Error updating template: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/email-templates/<template_id>', methods=['DELETE'])
@auth_required
def delete_email_template(template_id):
    """Delete email template"""
    try:
        user_email = session.get('user_email')
        success = postgres_db.delete_email_template(template_id, user_email)
        
        if success:
            return jsonify({'message': 'Template deleted successfully'})
        else:
            return jsonify({'error': 'Template not found or unauthorized'}), 404
    except Exception as e:
        print(f"[CUSTOMIZATION] Error deleting template: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/analytics', methods=['GET'])
@auth_required
def get_analytics():
    """Get analytics data with optional project filtering"""
    try:
        project_id = request.args.get('project_id')
        days = request.args.get('days', type=int)
        
        # Get current user
        user_id = session.get('user_email', 'unknown')
        
        # Get analytics from PostgreSQL
        if not days:
            days = 30
        
        analytics_data = postgres_db.get_analytics(user_id, days)
        if not analytics_data:  # Safety check
            analytics_data = {
                'total_campaigns': 0,
                'success_rate': 0,
                'outreach_types': {'founder': 0, 'investor': 0, 'any': 0},
                'projects': {},
                'recent_tasks': []
            }
        
        # Get recent tasks
        tasks = postgres_db.get_user_tasks(user_id, limit=10)
        if not tasks:  # Safety check
            tasks = []
        
        # Filter by project if specified
        if project_id and project_id != 'all':
            tasks = [t for t in tasks if t and str(t.get('project_id')) == str(project_id)]
            
            # Recalculate stats for filtered project
            filtered_analytics = {
                'total_campaigns': len(tasks),
                'success_rate': 100,
                'outreach_types': {},
                'projects': {}
            }
            
            for task in tasks:
                otype = task.get('outreach_type', 'unknown')
                filtered_analytics['outreach_types'][otype] = filtered_analytics['outreach_types'].get(otype, 0) + 1
            
            # Get project info
            project = postgres_db.get_project(project_id)
            if project:
                filtered_analytics['projects'][project['name']] = project.get('task_count', 0)
            
            analytics_data = filtered_analytics
        
        # Format recent tasks for frontend
        formatted_tasks = []
        for task in tasks[:10]:
            formatted_tasks.append({
                'task_id': task.get('task_id', task.get('id', '')),
                'outreach_type': task.get('outreach_type', 'unknown'),
                'project_name': 'General',  # Will be enhanced later
                'company_name': task.get('company_name', ''),
                'investor_name': task.get('investor_name', ''),
                'person_name': task.get('person_name', ''),
                'created_at': task.get('created_at', ''),
                'project_id': task.get('project_id', '1')
            })
        
        # Get project names for tasks
        projects = postgres_db.get_user_projects(user_id)
        if not projects:
            projects = []
        project_map = {str(p.get('id', '')): p.get('name', 'General') for p in projects if p}
        for task in formatted_tasks:
            task['project_name'] = project_map.get(str(task['project_id']), 'General')
        
        analytics_data['recent_tasks'] = formatted_tasks
        
        print(f"[ANALYTICS] Returning data for user {user_id}: total_campaigns={analytics_data.get('total_campaigns', 0)}")
        return jsonify(analytics_data)
        
    except Exception as e:
        print(f"[ERROR] Error getting analytics: {e}")
        import traceback
        print(f"[ERROR] Full traceback: {traceback.format_exc()}")
        return jsonify({'error': 'Failed to load analytics'}), 500

# Project management routes
@app.route('/api/projects', methods=['GET'])
@auth_required
def get_projects():
    """Get all projects"""
    projects = get_all_projects()
    # Add task count to each project
    for project in projects:
        project['task_count'] = get_project_task_count(project['id'])
    return jsonify(projects)

@app.route('/api/projects', methods=['POST'])
@auth_required
def create_new_project():
    """Create a new project"""
    data = request.get_json()
    if not data or 'name' not in data:
        return jsonify({'error': 'Project name is required'}), 400
    
    name = data['name'].strip()
    description = data.get('description', '').strip()
    color = data.get('color', '#4FD1C5')
    
    if not name:
        return jsonify({'error': 'Project name cannot be empty'}), 400
    
    project_id = create_project(name, description, color)
    if project_id is None:
        return jsonify({'error': 'Project name already exists'}), 400
    
    return jsonify({
        'id': project_id,
        'name': name,
        'description': description,
        'color': color,
        'task_count': 0
    }), 201

@app.route('/api/projects/<int:project_id>', methods=['DELETE'])
@auth_required
def delete_project_route(project_id):
    """Delete a project"""
    if project_id == 1:
        return jsonify({'error': 'Cannot delete the General project'}), 400
    
    success = delete_project(project_id)
    if success:
        return jsonify({'message': 'Project deleted successfully'})
    else:
        return jsonify({'error': 'Failed to delete project'}), 500

@app.route('/api/projects/<int:project_id>/tasks', methods=['GET'])
@auth_required
def get_project_tasks(project_id):
    """Get all tasks for a specific project"""
    try:
        # Get current user from session
        user_id = session.get('user_email', 'unknown')
        
        # Get tasks from PostgreSQL based on project_id
        from postgres_db import get_connection, return_connection
        
        if not postgres_db.pg_pool:
            return jsonify([])  # Return empty array if PostgreSQL not available
            
        conn = get_connection()
        try:
            with conn.cursor(cursor_factory=postgres_db.RealDictCursor) as cursor:
                # Get tasks for the specific project
                if str(project_id) == 'all':
                    cursor.execute("""
                        SELECT t.*, p.name as project_name 
                        FROM tasks t 
                        LEFT JOIN projects p ON t.project_id = p.project_id 
                        WHERE t.user_id = %s 
                        ORDER BY t.created_at DESC 
                        LIMIT 100
                    """, (user_id,))
                else:
                    cursor.execute("""
                        SELECT t.*, p.name as project_name 
                        FROM tasks t 
                        LEFT JOIN projects p ON t.project_id = p.project_id 
                        WHERE t.user_id = %s AND t.project_id = %s 
                        ORDER BY t.created_at DESC 
                        LIMIT 100
                    """, (user_id, str(project_id)))
                
                results = cursor.fetchall()
                tasks = []
                for row in results:
                    task_data = dict(row['task_data']) if row['task_data'] else {}
                    task = {
                        'task_id': row['task_id'],
                        'project_id': row['project_id'],
                        'outreach_type': row['outreach_type'],
                        'status': row['status'],
                        'project_name': row['project_name'] or 'General',
                        'created_at': row['created_at'].isoformat() if row['created_at'] else None,
                        'company_name': task_data.get('company_name'),
                        'investor_name': task_data.get('investor_name'),
                        'fund_name': task_data.get('fund_name'),
                        'person_name': task_data.get('person_name'),
                        'organization': task_data.get('organization'),
                        'email_content': task_data.get('email_content'),
                        'file_id': task_data.get('file_id')
                    }
                    tasks.append(task)
                
                return jsonify(tasks)
        finally:
            return_connection(conn)
            
    except Exception as e:
        print(f"[ERROR] Error getting project tasks: {e}")
        return jsonify({'error': 'Failed to load project tasks'}), 500

@app.route('/login-page')
def login_page():
    """Dedicated login page route"""
    return render_template('login.html')

@app.route('/logo')
def logo():
    """Serve the Westbridge logo"""
    logo_path = Path('wblogo.png')
    if logo_path.exists():
        return send_file(str(logo_path), mimetype='image/png')
    else:
        return '', 404

@app.route('/convert', methods=['POST'])
@auth_required
def convert():
    try:
        data = request.get_json()
        outreach_type = data.get('outreach_type', 'founder')
        recipient_email = data.get('recipient_email')  # Email of the person to send to
        cc_recipients = data.get('cc_recipients')  # CC email addresses
        project_id = data.get('project_id', 1)  # Default to General project
        
        # Get logged-in user's email and name from session
        user_email = session.get('user_email')
        user_name = session.get('user_name', 'Westbridge Capital Team')
        
        # Ensure project_id is valid and not None
        if not project_id or project_id == 'undefined' or project_id == 'null':
            project_id = 1
        
        if not recipient_email:
            return jsonify({'error': 'Recipient email address is required'}), 400
        
        # if not user_email:
        #     return jsonify({'error': 'User not authenticated'}), 401
        
        # Generate unique task ID
        task_id = str(uuid.uuid4())
        
        # Create progress queue
        progress_queues[task_id] = queue.Queue()
        
        if outreach_type == 'investor':
            # Handle investor research mode
            investor_name = data.get('investor_name')
            fund_name = data.get('fund_name')
            outreach_context = data.get('outreach_context')
            
            if not investor_name:
                return jsonify({'error': 'Investor name is required'}), 400
            if not fund_name:
                return jsonify({'error': 'Fund name is required'}), 400
            if not outreach_context:
                return jsonify({'error': 'Outreach context is required'}), 400
            
            # Start investor research in background thread
            thread = threading.Thread(
                target=run_investor_research_with_progress,
                args=(investor_name, fund_name, outreach_context, task_id, user_email, project_id, recipient_email, user_name, cc_recipients)
            )
            thread.start()
            
        elif outreach_type == 'any':
            # Handle any outreach mode
            person_name = data.get('person_name')
            organization = data.get('organization')
            outreach_request = data.get('outreach_request')
            
            if not person_name:
                return jsonify({'error': 'Person name is required'}), 400
            if not organization:
                return jsonify({'error': 'Organization is required'}), 400
            if not outreach_request:
                return jsonify({'error': 'Outreach request is required'}), 400
            
            # Start any outreach email generation in background thread
            thread = threading.Thread(
                target=run_any_outreach_with_progress,
                args=(person_name, organization, outreach_request, task_id, user_email, project_id, recipient_email, user_name, cc_recipients)
            )
            thread.start()
            
        else:
            # Handle founder outreach mode (existing logic)
            url = data.get('url')
            expert_info = data.get('expert_info', '').strip() if data.get('has_expert_info') else None
            
            if not url:
                return jsonify({'error': 'URL is required for founder outreach'}), 400
            
            output_path = app.config['TEMP_DIR'] / f"{task_id}.pptx"
            
            # Start conversion in background thread
            thread = threading.Thread(
                target=run_conversion_with_progress,
                args=(url, output_path, task_id, user_email, expert_info, project_id, recipient_email, user_name, cc_recipients)
            )
            thread.start()
        
        return jsonify({
            'success': True,
            'task_id': task_id
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/progress/<task_id>')
@auth_required
def progress(task_id):
    """Server-sent events endpoint for progress updates"""
    def generate():
        if task_id not in progress_queues:
            yield f"data: {json.dumps({'status': 'error', 'message': 'Invalid task ID'})}\n\n"
            return
        
        progress_queue = progress_queues[task_id]
        
        while True:
            try:
                # Get progress update
                progress = progress_queue.get(timeout=30)
                
                if progress is None:
                    # End of updates
                    del progress_queues[task_id]
                    break
                
                yield f"data: {json.dumps(progress)}\n\n"
                
            except queue.Empty:
                # Send heartbeat
                yield f"data: {json.dumps({'status': 'heartbeat'})}\n\n"
    
    return Response(generate(), mimetype="text/event-stream")

@app.route('/send-email', methods=['POST'])
@auth_required
def send_email():
    try:
        data = request.get_json()
        task_id = data.get('task_id')
        recipient = data.get('recipient')
        subject = data.get('subject')
        content = data.get('content')
        cc_recipients = data.get('cc_recipients')  # New CC field
        
        if not all([task_id, recipient, subject, content]):
            return jsonify({'error': 'Missing required fields'}), 400
        
        # Get generated content (from memory or database)
        task_data = get_task_metadata(task_id)
        if not task_data:
            return jsonify({'error': 'Invalid task ID or content not found'}), 404
        
        # Use PDF if available, otherwise fall back to PPTX
        attachment_path = task_data.get('pdf_path') or task_data.get('pptx_path')
        
        # Debug logging for attachment path
        print(f"[EMAIL] Task ID: {task_id}")
        print(f"[EMAIL] PDF path from DB: {task_data.get('pdf_path')}")
        print(f"[EMAIL] PPTX path from DB: {task_data.get('pptx_path')}")
        print(f"[EMAIL] Selected attachment path: {attachment_path}")
        
        # Check if file actually exists, with enhanced fallback logic
        if attachment_path:
            if os.path.exists(attachment_path):
                print(f"[EMAIL] Attachment file exists at: {attachment_path}")
                print(f"[EMAIL] File size: {os.path.getsize(attachment_path)} bytes")
            else:
                print(f"[EMAIL] WARNING: Attachment file NOT FOUND at: {attachment_path}")
                attachment_path = None
        
        # If no attachment path or file doesn't exist, try to find files using file_id
        if not attachment_path:
            file_id = task_data.get('file_id')
            print(f"[EMAIL] Attempting to find files using file_id: {file_id}")
            
            if file_id:
                # First try database
                db_file = postgres_db.get_generated_file(file_id, user_email)
                if db_file:
                    print(f"[EMAIL] Found file in database: {db_file['file_type']}")
                    # Create temporary file for email attachment
                    import tempfile
                    temp_fd, temp_path = tempfile.mkstemp(suffix=f".{db_file['file_type']}")
                    try:
                        with os.fdopen(temp_fd, 'wb') as tmp_file:
                            tmp_file.write(db_file['file_data'])
                        attachment_path = temp_path
                        print(f"[EMAIL] Created temporary file from database: {attachment_path}")
                        print(f"[EMAIL] File size: {len(db_file['file_data'])} bytes")
                    except Exception as e:
                        print(f"[EMAIL] Error creating temporary file from database: {e}")
                        os.close(temp_fd)
                        attachment_path = None
                else:
                    # Fallback to filesystem
                    print(f"[EMAIL] File not found in database, checking filesystem")
                    pdf_temp_path = app.config['TEMP_DIR'] / f"{file_id}.pdf"
                    pptx_temp_path = app.config['TEMP_DIR'] / f"{file_id}.pptx"
                    
                    print(f"[EMAIL] Checking PDF path: {pdf_temp_path}")
                    print(f"[EMAIL] Checking PPTX path: {pptx_temp_path}")
                    
                    if pdf_temp_path.exists():
                        attachment_path = str(pdf_temp_path)
                        print(f"[EMAIL] Found PDF in TEMP_DIR: {attachment_path}")
                        print(f"[EMAIL] PDF file size: {pdf_temp_path.stat().st_size} bytes")
                    elif pptx_temp_path.exists():
                        attachment_path = str(pptx_temp_path)
                        print(f"[EMAIL] Found PPTX in TEMP_DIR: {attachment_path}")
                        print(f"[EMAIL] PPTX file size: {pptx_temp_path.stat().st_size} bytes")
                    else:
                        print(f"[EMAIL] ERROR: No files found in TEMP_DIR for file_id: {file_id}")
                        print(f"[EMAIL] TEMP_DIR contents: {list(app.config['TEMP_DIR'].iterdir()) if app.config['TEMP_DIR'].exists() else 'TEMP_DIR does not exist'}")
            else:
                print(f"[EMAIL] ERROR: No file_id found in task_data")
                print(f"[EMAIL] Available task_data keys: {list(task_data.keys())}")
        
        # Get appropriate name based on outreach type
        outreach_type = task_data.get('outreach_type', 'founder')
        if outreach_type == 'investor':
            # For investor outreach, use investor name and fund name
            display_name = task_data.get('investor_name', 'Investor')
            if task_data.get('fund_name'):
                display_name += f" ({task_data['fund_name']})"
        elif outreach_type == 'any':
            # For any outreach, use person name and organization
            display_name = task_data.get('person_name', 'Contact')
            if task_data.get('organization'):
                display_name += f" at {task_data['organization']}"
        else:
            # For founder outreach, use company name
            display_name = task_data.get('company_name', 'Company')
        
        # Use the subject passed from frontend
        email_body = content
        
        # Generate attachment filename with company name
        attachment_name = None
        if attachment_path:
            print(f"[EMAIL] Final attachment check before send: {attachment_path}")
            print(f"[EMAIL] File exists: {os.path.exists(attachment_path)}")
            
            # Get company name for filename
            company_name = task_data.get('company_name')
            if company_name:
                # Determine file extension
                file_extension = Path(attachment_path).suffix.lower()
                if file_extension == '.pdf':
                    attachment_name = f"{company_name} and Westbridge Capital.pdf"
                elif file_extension == '.pptx':
                    attachment_name = f"{company_name} and Westbridge Capital.pptx"
                else:
                    attachment_name = f"{company_name} and Westbridge Capital{file_extension}"
                
                print(f"[EMAIL] Generated attachment name: {attachment_name}")
            else:
                print(f"[EMAIL] No company name found, using original filename")
        
        # Send email
        success, message = send_graph_email(
            recipient, 
            subject, 
            email_body, 
            attachment_path,
            attachment_name,
            cc_recipients
        )
        
        # Clean up temporary file if it was created from database
        if attachment_path and attachment_path.startswith('/tmp/'):
            try:
                os.unlink(attachment_path)
                print(f"[EMAIL] Cleaned up temporary file: {attachment_path}")
            except Exception as cleanup_error:
                print(f"[EMAIL] Warning: Failed to clean up temporary file {attachment_path}: {cleanup_error}")
        
        if success:
            return jsonify({'success': True, 'message': message})
        else:
            return jsonify({'error': message}), 500
            
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/download-pptx/<file_id>')
@auth_required
def download_pptx(file_id):
    try:
        # Validate file_id
        if not file_id or '..' in file_id or '/' in file_id:
            return jsonify({'error': 'Invalid file ID'}), 400
        
        user_email = session.get('user_email')
        
        # First try to get file from database
        db_file = postgres_db.get_generated_file(file_id, user_email)
        if db_file and db_file['file_type'] == 'pptx':
            print(f"[FILE STORAGE] Serving PPTX from database for {file_id}")
            
            # Use company name from database if available
            if db_file.get('company_name'):
                download_name = f"{db_file['company_name']}_and_Westbridge_Capital.pptx"
            else:
                download_name = db_file.get('filename', f"westbridge_outreach_deck_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
            
            # Create a temporary file-like object
            from io import BytesIO
            file_obj = BytesIO(db_file['file_data'])
            file_obj.seek(0)
            
            return send_file(
                file_obj,
                as_attachment=True,
                download_name=download_name,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
        
        # Fallback to filesystem
        print(f"[FILE STORAGE] File not found in database, checking filesystem for {file_id}")
        file_path = app.config['TEMP_DIR'] / f"{file_id}.pptx"
        
        if not file_path.exists():
            return jsonify({'error': 'File not found'}), 404
        
        # Try to extract company name from the PowerPoint file
        company_name = extract_company_name_from_pptx(str(file_path))
        
        # Generate filename using company name if available
        if company_name:
            download_name = f"{company_name}_and_Westbridge_Capital.pptx"
        else:
            # Fallback to timestamp-based naming
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            download_name = f"westbridge_outreach_deck_{timestamp}.pptx"
        
        return send_file(
            str(file_path),
            as_attachment=True,
            download_name=download_name,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/view-pdf/<file_id>')
@auth_required
def view_pdf(file_id):
    try:
        # Validate file_id
        if not file_id or '..' in file_id or '/' in file_id:
            return jsonify({'error': 'Invalid file ID'}), 400
        
        user_email = session.get('user_email')
        
        # First try to get file from database
        db_file = postgres_db.get_generated_file(file_id, user_email)
        if db_file and db_file['file_type'] == 'pdf':
            print(f"[FILE STORAGE] Serving PDF from database for {file_id}")
            
            # Create a temporary file-like object
            from io import BytesIO
            file_obj = BytesIO(db_file['file_data'])
            file_obj.seek(0)
            
            return send_file(
                file_obj,
                mimetype='application/pdf'
            )
        
        # Fallback to filesystem
        print(f"[FILE STORAGE] PDF not found in database, checking filesystem for {file_id}")
        pdf_path = app.config['TEMP_DIR'] / f"{file_id}.pdf"
        
        if not pdf_path.exists():
            return jsonify({'error': 'PDF not found'}), 404
        
        return send_file(
            str(pdf_path),
            mimetype='application/pdf'
        )
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/download-pdf/<file_id>')
@auth_required
def download_pdf(file_id):
    try:
        # Validate file_id
        if not file_id or '..' in file_id or '/' in file_id:
            return jsonify({'error': 'Invalid file ID'}), 400
        
        user_email = session.get('user_email')
        
        # First try to get file from database
        db_file = postgres_db.get_generated_file(file_id, user_email)
        if db_file and db_file['file_type'] == 'pdf':
            print(f"[FILE STORAGE] Serving PDF download from database for {file_id}")
            
            # Use company name from database if available
            if db_file.get('company_name'):
                download_name = f"{db_file['company_name']}_and_Westbridge_Capital.pdf"
            else:
                download_name = db_file.get('filename', f"westbridge_outreach_deck_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf")
            
            # Create a temporary file-like object
            from io import BytesIO
            file_obj = BytesIO(db_file['file_data'])
            file_obj.seek(0)
            
            return send_file(
                file_obj,
                as_attachment=True,
                download_name=download_name,
                mimetype='application/pdf'
            )
        
        # Fallback to filesystem
        print(f"[FILE STORAGE] PDF not found in database, checking filesystem for {file_id}")
        pdf_path = app.config['TEMP_DIR'] / f"{file_id}.pdf"
        
        if not pdf_path.exists():
            return jsonify({'error': 'PDF not found'}), 404
        
        # Try to extract company name from the corresponding PowerPoint file
        pptx_path = app.config['TEMP_DIR'] / f"{file_id}.pptx"
        company_name = None
        if pptx_path.exists():
            company_name = extract_company_name_from_pptx(str(pptx_path))
        
        # Generate filename using company name if available
        if company_name:
            download_name = f"{company_name}_and_Westbridge_Capital.pdf"
        else:
            # Fallback to timestamp-based naming
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            download_name = f"westbridge_outreach_deck_{timestamp}.pdf"
        
        return send_file(
            str(pdf_path),
            as_attachment=True,
            download_name=download_name,
            mimetype='application/pdf'
        )
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/history')
@auth_required
def get_generation_history():
    """Get all generation history from database for frontend"""
    try:
        user_email = session.get('user_email')
        #if not user_email:
        #    return jsonify({'error': 'User not authenticated'}), 401
        
        # Get tasks from PostgreSQL database
        tasks = postgres_db.get_user_tasks(user_email)
        
        # Format for frontend consumption
        history_items = []
        for task in tasks:
            # Extract task data safely
            task_data = task.get('task_data', {}) if task.get('task_data') else {}
            
            item = {
                'task_id': task.get('task_id'),
                'taskId': task.get('task_id'),  # For backward compatibility
                'outreach_type': task.get('outreach_type', 'founder'),
                'outreachType': task.get('outreach_type', 'founder'),  # For backward compatibility
                'status': task.get('status', 'completed'),
                'created_at': task.get('created_at'),
                'timestamp': task.get('created_at'),  # For backward compatibility
                'project_id': task.get('project_id'),
                'task_data': task_data,
                # Extract commonly used fields for easier access
                'url': task_data.get('url'),
                'investorName': task_data.get('investor_name'),
                'fundName': task_data.get('fund_name'),
                'personName': task_data.get('person_name'),
                'organization': task_data.get('organization'),
                'sector': task_data.get('sector'),
                'recipientEmail': task_data.get('recipient_email'),
                'emailContent': task_data.get('email_content'),
                'investorInsights': task_data.get('investor_insights'),
                'fileId': task_data.get('file_id')
            }
            history_items.append(item)
        
        return jsonify({
            'success': True,
            'history': history_items
        })
        
    except Exception as e:
        print(f"[HISTORY] Error getting history: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/debug-session')
def debug_session():
    """Debug endpoint to check session status"""
    session_info = {
        'user_authenticated': session.get('user_authenticated'),
        'user_email': session.get('user_email'),
        'user_name': session.get('user_name'),
        'session_keys': list(session.keys()),
        'session_permanent': session.permanent,
        'debug_mode': app.debug
    }
    return jsonify(session_info)

@app.route('/test-login')
def test_login():
    """Test endpoint to simulate login for development"""
    if app.debug:
        session['user_authenticated'] = True
        session['user_email'] = 'test@example.com'
        session['user_name'] = 'Test User'
        session.permanent = True
        return jsonify({'message': 'Test login successful', 'redirect': '/'})
    else:
        return jsonify({'error': 'Only available in debug mode'}), 403

@app.route('/test-notifications')
def test_notifications():
    """Test page for notifications"""
    if app.debug:
        from flask import send_file
        return send_file('test-notifications.html')
    else:
        return jsonify({'error': 'Only available in debug mode'}), 403


@app.route('/api/debug-analytics', methods=['GET'])
@auth_required
def debug_analytics():
    """Debug analytics data generation"""
    try:
        user_id = session.get('user_email', 'unknown')
        print(f"[DEBUG] User ID: {user_id}")
        
        # Check if we can get user tasks
        tasks = postgres_db.get_user_tasks(user_id, limit=5)
        print(f"[DEBUG] User tasks: {tasks}")
        
        # Check if we can get analytics
        analytics = postgres_db.get_analytics(user_id, 30)
        print(f"[DEBUG] Analytics: {analytics}")
        
        return jsonify({
            'user_id': user_id,
            'tasks_count': len(tasks) if tasks else 0,
            'analytics': analytics,
            'postgres_available': postgres_db.pg_pool is not None
        })
    except Exception as e:
        import traceback
        print(f"[DEBUG] Error: {e}")
        print(f"[DEBUG] Traceback: {traceback.format_exc()}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/postgres-health', methods=['GET'])
@auth_required
def postgres_health():
    """Check PostgreSQL connection health"""
    try:
        health = postgres_db.health_check()
        print(f"[HEALTH] PostgreSQL health check: {health}")
        return jsonify(health)
    except Exception as e:
        print(f"[HEALTH] PostgreSQL health check error: {e}")
        return jsonify({'error': str(e), 'status': 'error'}), 500

@app.route('/api/batch/upload', methods=['POST'])
@auth_required
def upload_batch_csv():
    """Upload and validate CSV file for batch processing"""
    try:
        if 'csv_file' not in request.files:
            return jsonify({'error': 'No CSV file provided'}), 400
        
        file = request.files['csv_file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not file.filename.lower().endswith('.csv'):
            return jsonify({'error': 'File must be a CSV'}), 400
        
        # Read CSV content
        stream = io.StringIO(file.stream.read().decode("UTF8"), newline=None)
        csv_input = csv.reader(stream)
        
        rows = []
        errors = []
        row_count = 0
        
        # Parse CSV rows
        for i, row in enumerate(csv_input):
            row_count += 1
            
            # Skip empty rows
            if not row or all(cell.strip() == '' for cell in row):
                continue
                
            # Check row limit
            if len(rows) >= 100:
                errors.append(f"CSV exceeds maximum of 100 rows")
                break
            
            # Validate row format
            if len(row) < 2:
                errors.append(f"Row {i+1}: Must have at least 2 columns (Company, Email)")
                continue
            
            company = row[0].strip()
            email = row[1].strip()
            
            # Validate company name
            if not company:
                errors.append(f"Row {i+1}: Company name cannot be empty")
                continue
            
            # Validate email format
            if not email:
                errors.append(f"Row {i+1}: Email cannot be empty")
                continue
            
            email_pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
            if not re.match(email_pattern, email):
                errors.append(f"Row {i+1}: Invalid email format: {email}")
                continue
            
            rows.append({
                'row_number': i + 1,
                'company': company,
                'email': email,
                'status': 'valid'
            })
        
        # Check minimum rows
        if len(rows) == 0:
            return jsonify({'error': 'No valid rows found in CSV'}), 400
        
        return jsonify({
            'success': True,
            'rows': rows,
            'total_rows': len(rows),
            'errors': errors,
            'has_errors': len(errors) > 0
        })
        
    except Exception as e:
        print(f"[CSV] Error parsing CSV: {e}")
        return jsonify({'error': f'Error parsing CSV: {str(e)}'}), 500

@app.route('/api/batch/submit', methods=['POST'])
@auth_required
def submit_batch_job():
    """Submit batch processing job to Anthropic Batch API"""
    print("[BATCH] ========== BATCH SUBMIT ENDPOINT CALLED ==========")
    print(f"[BATCH] Method: {request.method}")
    print(f"[BATCH] URL: {request.url}")
    print(f"[BATCH] Headers: {dict(request.headers)}")
    print("[BATCH] /api/batch/submit endpoint called")
    try:
        print("[BATCH] Getting request data...")
        data = request.get_json()
        print(f"[BATCH] Request data received: batch_type={data.get('batch_type') if data else None}, rows count={len(data.get('rows', [])) if data else 0}")
        
        if not data:
            print("[BATCH] ERROR: No data provided")
            return jsonify({'error': 'No data provided'}), 400
        
        batch_type = data.get('batch_type')
        rows = data.get('rows', [])
        sector = data.get('sector')
        
        print(f"[BATCH] Parsed data - batch_type: {batch_type}, rows: {len(rows)}, sector: {sector}")
        
        if not batch_type or batch_type not in ['founder', 'investor', 'any']:
            print(f"[BATCH] ERROR: Invalid batch type: {batch_type}")
            return jsonify({'error': 'Invalid batch type'}), 400
        
        if not rows:
            print("[BATCH] ERROR: No rows provided")
            return jsonify({'error': 'No rows provided'}), 400
        
        if len(rows) > 100:
            print(f"[BATCH] ERROR: Too many rows: {len(rows)}")
            return jsonify({'error': 'Too many rows (max 100)'}), 400
        
        user_email = session.get('user_email')
        print(f"[BATCH] User email from session: {user_email}")
        #if not user_email:
        #    print("[BATCH] ERROR: User not authenticated")
        #    return jsonify({'error': 'User not authenticated'}), 401
        
        # Create batch job record in database
        batch_job_id = str(uuid.uuid4())
        
        # Store batch job in database
        try:
            postgres_db.create_batch_job(
                batch_id=batch_job_id,
                user_email=user_email,
                batch_type=batch_type,
                sector=sector,
                total_rows=len(rows),
                csv_data=json.dumps(rows),
                status='pending'
            )
        except Exception as db_error:
            print(f"[BATCH] Database error: {db_error}")
            return jsonify({'error': 'Failed to create batch job record'}), 500
        
        # Submit to Anthropic Batch API
        try:
            print(f"[BATCH] About to submit to Anthropic Batch API")
            print(f"[BATCH] Batch job ID: {batch_job_id}")
            print(f"[BATCH] Batch type: {batch_type}")
            print(f"[BATCH] Number of rows: {len(rows)}")
            print(f"[BATCH] Sector: {sector}")
            
            # Use synchronous version for gevent compatibility
            anthropic_batch_id = submit_to_anthropic_batch_sync(batch_job_id, batch_type, rows, sector, user_email)
            
            print(f"[BATCH] Anthropic batch ID returned: {anthropic_batch_id}")
            
            if anthropic_batch_id:
                # Update batch job with Anthropic batch ID
                print(f"[BATCH] Updating batch job status to 'submitted'")
                postgres_db.update_batch_job_status(
                    batch_job_id, 
                    'submitted', 
                    anthropic_batch_id=anthropic_batch_id
                )
            else:
                print(f"[BATCH] ERROR: No batch ID returned from Anthropic")
                postgres_db.update_batch_job_status(
                    batch_job_id, 
                    'failed', 
                    error_message='Failed to submit to Anthropic Batch API'
                )
                return jsonify({'error': 'Failed to submit to Anthropic Batch API'}), 500
                
        except Exception as batch_error:
            print(f"[BATCH] Exception during Anthropic API submission: {type(batch_error).__name__}")
            print(f"[BATCH] Exception message: {str(batch_error)}")
            import traceback
            print(f"[BATCH] Exception traceback: {traceback.format_exc()}")
            postgres_db.update_batch_job_status(
                batch_job_id, 
                'failed', 
                error_message=str(batch_error)
            )
            return jsonify({'error': f'Batch API error: {str(batch_error)}'}), 500
        return jsonify({
            'success': True,
            'batch_id': batch_job_id,
            'status': 'pending',
            'message': 'Batch job submitted successfully'
        })
        
    except Exception as e:
        print(f"[BATCH] Error in submit_batch_job endpoint: {type(e).__name__}")
        print(f"[BATCH] Error message: {str(e)}")
        import traceback
        print(f"[BATCH] Full traceback: {traceback.format_exc()}")
        return jsonify({'error': f'Error submitting batch job: {str(e)}'}), 500

@app.route('/api/batch/list', methods=['GET'])
@auth_required
def list_batch_jobs():
    """Get all batch jobs for the current user"""
    try:
        user_email = session.get('user_email')
        #if not user_email:
        #    return jsonify({'error': 'User not authenticated'}), 401
        
        batch_jobs = postgres_db.get_user_batch_jobs(user_email)
        
        return jsonify({
            'success': True,
            'batch_jobs': batch_jobs
        })
        
    except Exception as e:
        print(f"[BATCH] Error listing batch jobs: {e}")
        return jsonify({'error': f'Error listing batch jobs: {str(e)}'}), 500

@app.route('/api/batch/notifications', methods=['GET'])
@auth_required  
def get_batch_notifications():
    """Get recent batch completion notifications for the user"""
    try:
        user_email = session.get('user_email')
        #if not user_email:
        #    return jsonify({'error': 'User not authenticated'}), 401
        
        # Get recently completed batch jobs
        conn = postgres_db.get_connection()
        try:
            with conn.cursor(cursor_factory=postgres_db.RealDictCursor) as cursor:
                cursor.execute("""
                    SELECT batch_id, batch_type, total_rows, completed_rows, 
                           status, completed_at, created_at
                    FROM batch_jobs 
                    WHERE user_id = %s 
                    AND status IN ('completed', 'failed')
                    AND completed_at > NOW() - INTERVAL '24 hours'
                    ORDER BY completed_at DESC
                    LIMIT 10
                """, (user_email,))
                
                notifications = []
                for job in cursor.fetchall():
                    notifications.append({
                        'id': job['batch_id'],
                        'type': 'batch_completed' if job['status'] == 'completed' else 'batch_failed',
                        'title': f"Batch {job['batch_type']} {'Completed' if job['status'] == 'completed' else 'Failed'}",
                        'message': f"Processed {job['completed_rows'] or 0}/{job['total_rows']} entries",
                        'timestamp': job['completed_at'].isoformat() if job['completed_at'] else None,
                        'status': job['status']
                    })
                
                return jsonify({
                    'success': True,
                    'notifications': notifications
                })
                
        finally:
            postgres_db.return_connection(conn)
            
    except Exception as e:
        print(f"[BATCH] Error getting notifications: {e}")
        return jsonify({'error': f'Error getting notifications: {str(e)}'}), 500

@app.route('/api/batch/status/<batch_id>', methods=['GET'])
@auth_required
def get_batch_status(batch_id):
    """Get status of a batch processing job"""
    try:
        user_email = session.get('user_email')
        batch_job = postgres_db.get_batch_job(batch_id, user_email)
        
        if not batch_job:
            return jsonify({'error': 'Batch job not found'}), 404
        
        return jsonify({
            'batch_id': batch_job['batch_id'],
            'status': batch_job['status'],
            'batch_type': batch_job['batch_type'],
            'total_rows': batch_job['total_rows'],
            'completed_rows': batch_job.get('completed_rows', 0),
            'created_at': batch_job['created_at'].isoformat() if batch_job['created_at'] else None,
            'completed_at': batch_job['completed_at'].isoformat() if batch_job.get('completed_at') else None,
            'error_message': batch_job.get('error_message')
        })
        
    except Exception as e:
        print(f"[BATCH] Error getting batch status: {e}")
        return jsonify({'error': f'Error getting batch status: {str(e)}'}), 500

@app.route('/api/batch/jobs', methods=['GET'])
@auth_required
def get_user_batch_jobs():
    """Get all batch jobs for the current user"""
    try:
        user_email = session.get('user_email')
        batch_jobs = postgres_db.get_user_batch_jobs(user_email)
        
        return jsonify({
            'batch_jobs': [
                {
                    'batch_id': job['batch_id'],
                    'status': job['status'],
                    'batch_type': job['batch_type'],
                    'total_rows': job['total_rows'],
                    'completed_rows': job.get('completed_rows', 0),
                    'created_at': job['created_at'].isoformat() if job['created_at'] else None,
                    'completed_at': job['completed_at'].isoformat() if job.get('completed_at') else None
                }
                for job in batch_jobs
            ]
        })
        
    except Exception as e:
        print(f"[BATCH] Error getting user batch jobs: {e}")
        return jsonify({'error': f'Error getting batch jobs: {str(e)}'}), 500

@app.route('/api/batch/results/<batch_id>', methods=['GET'])
@auth_required
def get_batch_results(batch_id):
    """Get results for a completed batch job"""
    try:
        user_email = session.get('user_email')
        batch_job = postgres_db.get_batch_job(batch_id, user_email)
        
        if not batch_job:
            return jsonify({'error': 'Batch job not found'}), 404
        
        if batch_job['status'] != 'completed':
            return jsonify({'error': 'Batch job not completed yet'}), 400
        
        # Parse results JSON
        results = json.loads(batch_job.get('results', '[]'))
        
        return jsonify({
            'batch_id': batch_id,
            'batch_type': batch_job['batch_type'],
            'total_rows': batch_job['total_rows'],
            'completed_rows': batch_job['completed_rows'],
            'results': results,
            'created_at': batch_job['created_at'].isoformat() if batch_job['created_at'] else None,
            'completed_at': batch_job['completed_at'].isoformat() if batch_job['completed_at'] else None
        })
        
    except Exception as e:
        print(f"[BATCH] Error getting batch results: {e}")
        return jsonify({'error': f'Error getting batch results: {str(e)}'}), 500


# ===== FINDER ROUTES ===== 

@app.route('/api/finder/analyze', methods=['POST'])
@auth_required
def analyze_finder_query():
    """Analyze the initial finder query"""
    try:
        data = request.get_json()
        if not data or not data.get('query'):
            return jsonify({'error': 'Query is required'}), 400
        
        query = data['query']
        lead_count = data.get('lead_count', 50)  # Default to 50 if not provided
        user_email = session.get('user_email')
        
        print(f"[FINDER] Analyzing query: {query[:100]}...")
        print(f"[FINDER] Requested lead count: {lead_count}")
        result = finder_service.process_initial_query(query, user_email, lead_count)
        
        return jsonify(result)
        
    except Exception as e:
        print(f"[FINDER] Error analyzing query: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/finder/clarification', methods=['POST'])
@auth_required
def submit_finder_clarification():
    """Submit clarification responses"""
    try:
        print("[FINDER] Received clarification request")
        data = request.get_json()
        print(f"[FINDER] Clarification data: {data}")
        
        if not data or not data.get('session_id') or not data.get('responses'):
            error_msg = 'Session ID and responses are required'
            print(f"[FINDER] Error: {error_msg}")
            return jsonify({'error': error_msg}), 400
        
        session_id = data['session_id']
        responses = data['responses']
        skip = data.get('skip', False)
        
        print(f"[FINDER] Processing clarification for session: {session_id}")
        print(f"[FINDER] Responses: {responses}")
        print(f"[FINDER] Skip: {skip}")
        
        result = finder_service.process_clarification_response(session_id, responses, skip)
        print(f"[FINDER] Clarification result: {result}")
        
        return jsonify(result)
        
    except Exception as e:
        print(f"[FINDER] Error processing clarification: {e}")
        import traceback
        print(f"[FINDER] Full traceback: {traceback.format_exc()}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/finder/search', methods=['POST'])
@auth_required
def execute_finder_search():
    """Execute the finder search"""
    try:
        print("[FINDER] Received search request")
        data = request.get_json()
        print(f"[FINDER] Request data: {data}")
        
        if not data or not data.get('session_id') or not data.get('orchestration'):
            error_msg = 'Session ID and orchestration data are required'
            print(f"[FINDER] Error: {error_msg}")
            return jsonify({'error': error_msg}), 400
        
        session_id = data['session_id']
        orchestration = data['orchestration']
        
        print(f"[FINDER] Executing search for session: {session_id}")
        print(f"[FINDER] Orchestration data: {orchestration}")
        
        result = finder_service.execute_search(session_id, orchestration)
        print(f"[FINDER] Search result: {result}")
        
        # Store CSV content in Flask session for download
        if result.get('csv_content'):
            session[f'finder_csv_{session_id}'] = result['csv_content']
            print(f"[FINDER] Stored CSV content in session")
        else:
            print(f"[FINDER] No CSV content in result to store")
        
        return jsonify(result)
        
    except Exception as e:
        print(f"[FINDER] Error executing search: {e}")
        import traceback
        print(f"[FINDER] Full traceback: {traceback.format_exc()}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/finder/download/<session_id>')
@auth_required
def download_finder_csv(session_id):
    """Download the CSV results from a finder search"""
    try:
        # For now, we'll store the CSV in the session or a temporary location
        # In a real implementation, you'd retrieve from database or cache
        csv_content = session.get(f'finder_csv_{session_id}')
        
        if not csv_content:
            return jsonify({'error': 'CSV not found'}), 404
        
        # Create response with CSV content
        response = Response(
            csv_content,
            mimetype='text/csv',
            headers={'Content-Disposition': f'attachment; filename=finder_results_{session_id[:8]}.csv'}
        )
        
        return response
        
    except Exception as e:
        print(f"[FINDER] Error downloading CSV: {e}")
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    # Environment validation is now handled in config.py
    print("\n[STARTUP] Internal Tooling - Outreach PowerPoint & Email Generator")
    print("[INFO] Open http://localhost:5100 in your browser")
    print("[EMAIL] Email sending:", "Enabled" if app.config.get('OUTLOOK_EMAIL') else "Disabled")
    
    # Check PostgreSQL health on startup
    postgres_health = postgres_db.health_check()
    if postgres_health and postgres_health.get('connected'):
        print(f"[POSTGRES] Connected - Version: {postgres_health['version']}, Records: {postgres_health['total_keys']}")
    else:
        error_msg = postgres_health.get('error') if postgres_health else 'Health check returned None'
        print(f"[POSTGRES] Not connected - {error_msg}")
    
    print()
    
    # Ensure leader election table exists
    ensure_leader_table()
    
    # Start background batch monitoring
    print("[BATCH] Starting background batch status monitor...")
    start_background_batch_monitor()
    
    app.run(debug=app.config.get('DEBUG', False), host='0.0.0.0', port=5100)