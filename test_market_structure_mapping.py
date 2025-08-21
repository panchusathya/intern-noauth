#!/usr/bin/env python3
"""
Test script to verify the new market structure 4-box layout mapping
"""

import sys
from pptx import Presentation
from website_to_ppt import update_market_structure_slide, remove_citations

def test_market_structure_mapping():
    """Test the new market structure slide mapping"""
    
    # Load the template
    template_path = "Unbound and WB Capital.pptx"
    
    try:
        prs = Presentation(template_path)
    except FileNotFoundError:
        print(f"Template file '{template_path}' not found!")
        return False
    
    # Get slide 3 (index 3, which is slide 4 in PowerPoint)
    if len(prs.slides) <= 3:
        print("Template doesn't have enough slides!")
        return False
    
    slide = prs.slides[3]  # Slide 4 (0-based index)
    
    # Test data matching the example from the image
    test_data = {
        "title": "Some Key Observations about Market Structure",
        "observations": [
            {
                "heading": "Memory Is Still A Huge Bottleneck",
                "content": "Over the last few years, compute power has scaled significantly faster than Memory and Bandwidth, leading to latency issues. As parameter size increases in ML models, the delay in data transfer between the CPU memory module to the tensor core which perform the matrix operations become significantly longer."
            },
            {
                "heading": "Energy Consumption", 
                "content": "Megatron-LM language model (smaller than GPT-3) trained by Nvidia over nine days consumed 2.6x the annual power consumption of avg. U.S. house. Issues with cooling AI chips are leading to server meltdowns at some CSPs. Cooling technology now in focus, data centers implementing various liquid cooling technology."
            },
            {
                "heading": "Profitability",
                "content": "AI compute demand scales linearly with developer headcount (i.e. when training models). At scale, GenAI companies could still face meaningful compute costs and fail to reach SaaS profitability levels. For example, when OpenAI was starting up in 2019-2020, they spent over 120M+ for GCP services."
            },
            {
                "heading": "Separation of Training and Inference",
                "content": "One way of addressing profitability issues is to separate inference from more intensives training workload and energy-hungry chips, as it requires significantly less compute. Whole slew of startups working to address the issues (ie. Groq, d-Matrix, FuriosaAI, NeuReality, SambaNova, and more)"
            }
        ]
    }
    
    print("Testing market structure slide mapping...")
    print(f"Slide has {len(slide.shapes)} shapes")
    
    # List all shapes for debugging
    print("\nAll shapes in slide:")
    for i, shape in enumerate(slide.shapes):
        shape_type = shape.shape_type
        has_text = shape.has_text_frame
        text_preview = shape.text[:30] + "..." if has_text and shape.text else "No text"
        print(f"  Shape {i}: {shape.name} (Type: {shape_type}, Text: {has_text}) - {text_preview}")
    
    # Test the update function
    try:
        update_market_structure_slide(slide, test_data)
        print("\n✅ Successfully updated market structure slide!")
        
        # Save test output
        output_path = "test_market_structure_output.pptx"
        prs.save(output_path)
        print(f"✅ Test output saved to: {output_path}")
        
        return True
        
    except Exception as e:
        print(f"\n❌ Error updating slide: {e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    success = test_market_structure_mapping()
    sys.exit(0 if success else 1)