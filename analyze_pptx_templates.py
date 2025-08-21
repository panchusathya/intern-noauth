#!/usr/bin/env python3
"""
Analyze PowerPoint templates to understand their structure and slide layouts
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_slide(slide, slide_index):
    """Analyze a single slide and return its structure"""
    print(f"\n{'='*80}")
    print(f"Slide {slide_index + 1}: {slide.slide_layout.name if hasattr(slide, 'slide_layout') else 'Unknown Layout'}")
    print(f"{'='*80}")
    
    # Get slide title if present
    if slide.shapes.title:
        print(f"Title: {slide.shapes.title.text}")
    
    print(f"\nTotal shapes: {len(slide.shapes)}")
    
    # Analyze each shape
    for shape_idx, shape in enumerate(slide.shapes):
        print(f"\nShape {shape_idx + 1}:")
        print(f"  Type: {shape.shape_type}")
        print(f"  Name: {shape.name}")
        
        # Position and size
        print(f"  Position: left={shape.left}, top={shape.top}")
        print(f"  Size: width={shape.width}, height={shape.height}")
        
        # Check if it's a text box
        if shape.has_text_frame:
            print(f"  Has text frame: Yes")
            print(f"  Text: {shape.text[:100]}..." if len(shape.text) > 100 else f"  Text: {shape.text}")
            print(f"  Paragraphs: {len(shape.text_frame.paragraphs)}")
            
            # Check text formatting
            if shape.text_frame.paragraphs:
                first_para = shape.text_frame.paragraphs[0]
                if first_para.runs:
                    first_run = first_para.runs[0]
                    if first_run.font.size:
                        print(f"  Font size: {first_run.font.size.pt}pt")
                    if first_run.font.name:
                        print(f"  Font name: {first_run.font.name}")
        
        # Check if it's a placeholder
        try:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                print(f"  Placeholder index: {shape.placeholder_format.idx}")
                print(f"  Placeholder type: {shape.placeholder_format.type}")
        except ValueError:
            # Shape is not a placeholder
            print(f"  Is placeholder: No")
        
        # Check if it's a group shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"  Group contains {len(shape.shapes)} shapes")
            
        # Check if it has fill
        if hasattr(shape, 'fill'):
            print(f"  Has fill: {shape.fill.type}")

def find_market_structure_slide(prs):
    """Find the 'Some Key Observations about Market Structure' slide"""
    for idx, slide in enumerate(prs.slides):
        if slide.shapes.title and "market structure" in slide.shapes.title.text.lower():
            return idx, slide
    return None, None

def analyze_template(template_path):
    """Analyze a PowerPoint template"""
    print(f"\nAnalyzing template: {template_path}")
    print("="*100)
    
    prs = Presentation(template_path)
    
    # Find and analyze the market structure slide specifically
    market_idx, market_slide = find_market_structure_slide(prs)
    
    if market_slide:
        print(f"\n{'*'*100}")
        print(f"FOUND MARKET STRUCTURE SLIDE AT INDEX {market_idx}")
        print(f"{'*'*100}")
        analyze_slide(market_slide, market_idx)
        
        # Extract the structure for mapping
        print(f"\n{'*'*100}")
        print("MAPPING STRUCTURE FOR MARKET OBSERVATIONS:")
        print(f"{'*'*100}")
        
        text_boxes = []
        for shape in market_slide.shapes:
            if shape.has_text_frame and not (shape == market_slide.shapes.title):
                text_boxes.append({
                    'name': shape.name,
                    'text': shape.text,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'has_fill': hasattr(shape, 'fill') and shape.fill.type is not None
                })
        
        # Sort by position (top to bottom, left to right)
        text_boxes.sort(key=lambda x: (x['top'], x['left']))
        
        print(f"\nFound {len(text_boxes)} text boxes (excluding title)")
        for idx, box in enumerate(text_boxes):
            print(f"\nText Box {idx + 1}:")
            print(f"  Name: {box['name']}")
            print(f"  Current text: {box['text'][:50]}..." if len(box['text']) > 50 else f"  Current text: {box['text']}")
            print(f"  Has background fill: {box['has_fill']}")
            
            # Determine if it's a heading or content box based on text patterns
            if any(keyword in box['text'].lower() for keyword in ['memory', 'energy', 'profitability', 'separation']):
                print(f"  Type: HEADING BOX")
            else:
                print(f"  Type: CONTENT BOX")
    else:
        print("\nMarket Structure slide not found. Analyzing all slides...")
        for idx, slide in enumerate(prs.slides):
            analyze_slide(slide, idx)

def main():
    """Main function"""
    # Check for template files in the current directory
    template_files = []
    
    # Look for Unbound and WB Capital templates
    for pattern in ['*Unbound*.pptx', '*WB*Capital*.pptx', '*westbridge*.pptx']:
        template_files.extend(Path('.').glob(pattern))
    
    if not template_files:
        print("No template files found. Please ensure Unbound and WB Capital PPTX files are in the current directory.")
        sys.exit(1)
    
    print(f"Found {len(template_files)} template(s):")
    for f in template_files:
        print(f"  - {f}")
    
    # Analyze each template
    for template_file in template_files:
        analyze_template(str(template_file))
        print("\n" + "="*100 + "\n")

if __name__ == "__main__":
    main()